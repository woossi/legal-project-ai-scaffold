#!/usr/bin/env python3
"""PPT 덱을 읽기 전용으로 검사하고 03_qa.json·03_qa.md를 작성한다."""

from __future__ import annotations

import argparse
import hashlib
import importlib.util
import json
import os
import re
import zipfile
from collections import Counter
from datetime import date
from decimal import Decimal, InvalidOperation
from pathlib import Path
from typing import Any

import jsonschema
from PIL import Image
from pptx import Presentation


SCRIPT = Path(__file__).resolve()
ROOT = next(parent for parent in SCRIPT.parents if (parent / "AGENTS.md").is_file())
SOURCE_SCHEMA = ROOT / ".claude/skills/ppt/ppt-source/contract/source.schema.json"
DECK_SCHEMA = ROOT / ".claude/skills/ppt/ppt-deck/contract/deck.schema.json"
BUILD_SCHEMA = ROOT / ".claude/skills/ppt/ppt-render/contract/build.schema.json"
QA_SCHEMA = ROOT / ".claude/skills/ppt/ppt-verify/contract/qa.schema.json"
QA_CONTRACT = ROOT / ".claude/skills/ppt/ppt-verify/contract/outputs.json"
MEASURE_SCRIPT = ROOT / ".claude/skills/ppt/ppt-source/scripts/measure.py"
MEASURE_SPEC = importlib.util.spec_from_file_location("ppt_source_measure", MEASURE_SCRIPT)
if MEASURE_SPEC is None or MEASURE_SPEC.loader is None:
    raise ImportError(f"실측 모듈을 불러올 수 없음: {MEASURE_SCRIPT}")
MEASURE_MODULE = importlib.util.module_from_spec(MEASURE_SPEC)
MEASURE_SPEC.loader.exec_module(MEASURE_MODULE)
measure_value = MEASURE_MODULE.measure_value

SOURCE_CHECKS = ("source_asset_hashes", "source_fact_integrity", "fact_value_check")
DECK_CHECKS = (
    "title_style",
    "slide_sequence",
    "cover_role",
    "bullet_hierarchy",
    "body_style",
    "text_overflow",
)
FACT_CHECKS = (
    "fact_refs_present",
    "display_value_binding",
    "fact_refs_resolve",
    "denominator_pair",
    "closing_on_limits",
)
PPTX_CHECKS = ("zip_integrity", "slide_count", "pptx_text_integrity")
BUILD_CHECKS = (
    "artifact_identity",
    "upstream_hash",
    "spec_hash",
    "previews_present",
    "artifact_paths",
    "font_fallback",
    "visual_check",
)
BLIND_SPOTS = [
    "PowerPoint 또는 Keynote의 네이티브 렌더 결과를 직접 비교하지 않음",
    "본문 문체 판정은 형태소 분석이 아닌 종결 패턴 검사에 제한됨",
    "길이 기반 오버플로 검사는 실제 글꼴 메트릭 겹침을 완전하게 판정하지 못함",
]

NUMBER_TOKEN_RE = re.compile(r"\d[\d,]*(?:\.\d+)?")
ISO_DATE_RE = re.compile(r"(?<!\d)\d{4}[-./]\d{1,2}[-./]\d{1,2}(?!\d)")
FILE_IDENTIFIER_RE = re.compile(
    r"(?<![A-Za-z0-9_.-])[\w.-]+\.(?:csv|html|json|jsonl|md|pdf|png|pptx|svg|ttl)"
    r"(?![A-Za-z0-9_.-])",
    re.IGNORECASE,
)
RATIO_RE = re.compile(r"[%％]|비율|대비|퍼센트")
QUESTION_END_RE = re.compile(r"[?？]\s*$")
QUESTION_FORM_RE = re.compile(
    r"(?:^(?:왜|무엇|어떻게|누가|언제|어디|얼마)\s|"
    r"(?:무엇인가|필요한가|가능한가|있는가|없는가|하는가|되는가|보이는가)\s*$)"
)
TASK_TITLE_RE = re.compile(r"^(?:Task|작업)\s+[1-9]\d*\s+—\s+\S.+$")
SENTENCE_END_RE = re.compile(
    r"(?:했다|한다|이다|있다|된다|됐다|였다|었다|았다|한다|다|지만|으며|면서|고)"
    r"[.!?。]?\s*$"
)
NOMINAL_END_RE = re.compile(r"(?:함|됨|임|음)[.!?。]?\s*$")
NOMINALIZED_TITLE_END_RE = re.compile(
    r"(?:했음|되었음|됐음|있었음|없었음|함|됨|임)[.!?。]?\s*$"
)


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _item(
    name: str,
    status: str,
    detail: str | None = None,
    *,
    slide_no: int | None = None,
    owner: str | None = None,
    threshold_overridden: bool = False,
) -> dict[str, Any]:
    if status == "fail" and (not detail or owner is None):
        raise ValueError(f"fail 항목은 detail과 owner가 필요함: {name}")
    return {
        "name": name,
        "status": status,
        "detail": detail,
        "slide_no": slide_no,
        "owner": owner,
        "threshold_overridden": threshold_overridden,
    }


def _append_unchecked(items: list[dict[str, Any]], names: tuple[str, ...], detail: str) -> None:
    for name in names:
        items.append(_item(name, "unchecked", detail))


def _load_and_validate_json(
    path: Path,
    schema_path: Path,
    name: str,
    owner: str,
    items: list[dict[str, Any]],
) -> tuple[dict[str, Any] | None, bool]:
    if not path.is_file():
        items.append(_item(name, "unchecked", "입력 파일 없음"))
        return None, False

    try:
        raw_value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, UnicodeDecodeError, json.JSONDecodeError) as error:
        items.append(_item(name, "fail", f"JSON 로드 실패: {error}", owner=owner))
        return None, False

    value = raw_value if isinstance(raw_value, dict) else None
    try:
        jsonschema.validate(raw_value, _load_json(schema_path))
    except jsonschema.ValidationError as error:
        location = "/".join(str(part) for part in error.absolute_path) or "<root>"
        items.append(_item(name, "fail", f"{location}: {error.message}", owner=owner))
        return value, False
    items.append(_item(name, "pass", "JSON Schema 통과"))
    return value, True


def _sentence_title(text: str) -> bool:
    stripped = text.strip()
    return bool(
        SENTENCE_END_RE.search(stripped)
        or NOMINALIZED_TITLE_END_RE.search(stripped)
        or QUESTION_END_RE.search(stripped)
        or QUESTION_FORM_RE.search(stripped)
    )


def _nominal_body(text: str) -> bool:
    stripped = text.strip()
    if NOMINAL_END_RE.search(stripped):
        return True
    return not _sentence_title(stripped)


def _source_asset_path(path_value: str) -> Path:
    path = Path(path_value)
    return path if path.is_absolute() else ROOT / path


def _append_title_checks(deck: dict[str, Any], items: list[dict[str, Any]]) -> None:
    violations: list[tuple[int | None, str]] = []
    if _sentence_title(deck.get("title", "")):
        violations.append((None, deck["title"]))
    for slide in deck.get("slides", []):
        title = slide.get("title", "")
        if _sentence_title(title):
            violations.append((slide.get("no"), slide["title"]))
        elif slide.get("role") == "task" and not TASK_TITLE_RE.fullmatch(title):
            violations.append((slide.get("no"), title))
        elif slide.get("role") == "cover" and title != deck.get("title"):
            violations.append((slide.get("no"), title))
    if not violations:
        items.append(_item("title_style", "pass", "덱·슬라이드 제목 전건 명사구"))
        return
    for slide_no, title in violations:
        items.append(
            _item(
                "title_style",
                "fail",
                f"서술형 제목: {title}",
                slide_no=slide_no,
                owner="ppt-writer",
            )
        )


def _append_slide_structure(deck: dict[str, Any], items: list[dict[str, Any]]) -> None:
    slides = deck.get("slides", [])
    numbers = [slide.get("no") for slide in slides]
    expected = list(range(1, len(slides) + 1))
    if numbers == expected:
        items.append(_item("slide_sequence", "pass", "슬라이드 번호 1부터 연속 증가"))
    else:
        items.append(
            _item(
                "slide_sequence",
                "fail",
                f"실제 번호 {numbers}, 예상 번호 {expected}",
                owner="ppt-writer",
            )
        )

    covers = [slide for slide in slides if slide.get("role") == "cover"]
    if len(covers) == 1 and covers[0].get("no") == 1:
        items.append(_item("cover_role", "pass", "cover 1개이며 첫 슬라이드"))
    else:
        items.append(
            _item(
                "cover_role",
                "fail",
                f"cover {len(covers)}개, 번호 {[slide.get('no') for slide in covers]}",
                owner="ppt-writer",
            )
        )


def _append_body_checks(deck: dict[str, Any], items: list[dict[str, Any]], threshold: float) -> None:
    bodies = [
        body
        for slide in deck.get("slides", [])
        for body in slide.get("body", [])
        if body.get("text", "").strip()
    ]
    if not bodies:
        items.append(_item("body_style", "warn", "본문 항목 없음", owner="ppt-writer"))
        return
    nominal = sum(1 for body in bodies if _nominal_body(body["text"]))
    ratio = nominal / len(bodies)
    if ratio < threshold:
        items.append(
            _item(
                "body_style",
                "warn",
                f"명사형 종결 {nominal}/{len(bodies)} ({ratio:.1%}), 기본 임계 {threshold:.0%} 미만",
                owner="ppt-writer",
            )
        )
    else:
        items.append(_item("body_style", "pass", f"명사형 종결 {nominal}/{len(bodies)} ({ratio:.1%})"))


def _append_bullet_hierarchy(deck: dict[str, Any], items: list[dict[str, Any]]) -> None:
    failures: list[tuple[int, str]] = []
    for slide in deck.get("slides", []):
        if slide.get("role") != "task":
            continue
        body = slide.get("body", [])
        for index, block in enumerate(body):
            if block.get("level", 1) != 1:
                continue
            has_child = index + 1 < len(body) and body[index + 1].get("level", 1) == 2
            if not has_child:
                failures.append(
                    (
                        slide["no"],
                        f"상위 불릿의 하위 근거 누락: {block.get('kind')} — {block.get('text')}",
                    )
                )

    if failures:
        for slide_no, detail in failures:
            items.append(
                _item(
                    "bullet_hierarchy",
                    "fail",
                    detail,
                    slide_no=slide_no,
                    owner="ppt-writer",
                )
            )
    else:
        items.append(_item("bullet_hierarchy", "pass", "task 상위 불릿 전건에 하위 근거 존재"))


def _decimal_value(value: Any) -> Decimal | None:
    if isinstance(value, bool) or value is None:
        return None
    try:
        return Decimal(str(value).replace(",", ""))
    except InvalidOperation:
        return None


def _display_matches_fact(display_value: Any, fact_value: Any) -> bool:
    display_number = _decimal_value(display_value)
    fact_number = _decimal_value(fact_value)
    if display_number is not None and fact_number is not None:
        return display_number == fact_number
    return type(display_value) is type(fact_value) and display_value == fact_value


def _fact_number_tokens(text: str) -> list[str]:
    """날짜와 파일 식별자를 제외한 발표 수치 토큰을 반환한다."""
    without_dates = ISO_DATE_RE.sub("", text)
    without_identifiers = FILE_IDENTIFIER_RE.sub("", without_dates)
    return NUMBER_TOKEN_RE.findall(without_identifiers)


def _append_fact_checks(source: dict[str, Any], deck: dict[str, Any], items: list[dict[str, Any]]) -> None:
    facts = {fact["id"]: fact for fact in source.get("facts", [])}
    missing_refs: list[tuple[int, str]] = []
    dangling: list[tuple[int, str]] = []
    denominator_errors: list[tuple[int, str]] = []
    binding_errors: list[tuple[int, str]] = []

    for slide in deck.get("slides", []):
        slide_no = slide["no"]
        blocks = list(slide.get("body", []))
        blocks.extend(
            {
                "text": visual.get("caption") or "",
                "fact_refs": visual.get("fact_refs", []),
                "claims": visual.get("claims", []),
            }
            for visual in slide.get("visuals", [])
        )
        for block in blocks:
            text = block.get("text", "")
            refs = block.get("fact_refs", [])
            claims = block.get("claims", [])
            tokens = _fact_number_tokens(text)
            if tokens and not refs:
                missing_refs.append((slide_no, text))
            for ref in refs:
                if ref not in facts:
                    dangling.append((slide_no, ref))

            claimed_numbers: list[Decimal] = []
            if tokens and not claims:
                binding_errors.append((slide_no, f"수치 claims 누락: {text}"))
            for claim in claims:
                ref = claim["fact_ref"]
                display_value = claim["display_value"]
                if ref not in refs:
                    binding_errors.append((slide_no, f"claim fact_ref가 fact_refs에 없음: {ref}"))
                if ref in facts and not _display_matches_fact(display_value, facts[ref].get("value")):
                    binding_errors.append(
                        (
                            slide_no,
                            f"표시값 {display_value!r}와 {ref} 값 {facts[ref].get('value')!r} 불일치",
                        )
                    )
                display_number = _decimal_value(display_value)
                if display_number is not None:
                    claimed_numbers.append(display_number)

            for token in tokens:
                token_number = _decimal_value(token)
                if token_number is not None and token_number not in claimed_numbers:
                    binding_errors.append((slide_no, f"표시 수치 {token}의 claim 누락: {text}"))

            if tokens and RATIO_RE.search(text):
                referenced = [facts[ref] for ref in refs if ref in facts]
                grounded = False
                for fact in referenced:
                    denominator_fact = fact.get("denominator_fact")
                    if fact.get("denominator") is not None:
                        grounded = True
                    elif denominator_fact and denominator_fact in refs:
                        grounded = True
                if not grounded:
                    denominator_errors.append((slide_no, text))

    if missing_refs:
        for slide_no, text in missing_refs:
            items.append(
                _item(
                    "fact_refs_present",
                    "fail",
                    f"수치 근거 누락: {text}",
                    slide_no=slide_no,
                    owner="ppt-writer",
                )
            )
    else:
        items.append(_item("fact_refs_present", "pass", "수치 포함 본문·캡션 전건에 fact_refs 존재"))

    if binding_errors:
        for slide_no, detail in binding_errors:
            items.append(
                _item(
                    "display_value_binding",
                    "fail",
                    detail,
                    slide_no=slide_no,
                    owner="ppt-writer",
                )
            )
    else:
        items.append(_item("display_value_binding", "pass", "표시 수치 전건 fact 값과 일치"))

    if dangling:
        for slide_no, ref in dangling:
            items.append(
                _item(
                    "fact_refs_resolve",
                    "fail",
                    f"존재하지 않는 fact id: {ref}",
                    slide_no=slide_no,
                    owner="ppt-curator",
                )
            )
    else:
        items.append(_item("fact_refs_resolve", "pass", "fact_refs 전건 해소"))

    if denominator_errors:
        for slide_no, text in denominator_errors:
            items.append(
                _item(
                    "denominator_pair",
                    "fail",
                    f"비율의 분모 근거 누락: {text}",
                    slide_no=slide_no,
                    owner="ppt-curator",
                )
            )
    else:
        items.append(_item("denominator_pair", "pass", "비율 항목 분모 근거 확인"))


def _append_source_asset_hashes(source: dict[str, Any], items: list[dict[str, Any]]) -> None:
    failures: list[str] = []
    for entry in source.get("sources", []):
        path = _source_asset_path(entry["path"])
        if not path.is_file():
            failures.append(f"{entry['id']}: 파일 없음 {entry['path']}")
        elif sha256_file(path) != entry["sha256"]:
            failures.append(f"{entry['id']}: SHA-256 불일치")
    if failures:
        for detail in failures:
            items.append(_item("source_asset_hashes", "fail", detail, owner="ppt-curator"))
    else:
        items.append(_item("source_asset_hashes", "pass", "sources 전건 SHA-256 일치"))


def _append_source_fact_integrity(source: dict[str, Any], items: list[dict[str, Any]]) -> None:
    sources = source.get("sources", [])
    facts = source.get("facts", [])
    source_ids = [entry.get("id") for entry in sources]
    fact_ids = [fact.get("id") for fact in facts]
    source_by_id = {entry.get("id"): entry for entry in sources}
    failures: list[str] = []

    if len(source_ids) != len(set(source_ids)):
        failures.append("sources[].id 중복")
    if len(fact_ids) != len(set(fact_ids)):
        failures.append("facts[].id 중복")
    for fact in facts:
        source_entry = source_by_id.get(fact.get("source_id"))
        if source_entry is None:
            failures.append(f"{fact.get('id')}: source_id 미해소 {fact.get('source_id')}")
        elif source_entry.get("contract_status") == "failed" and not fact.get("caveat"):
            failures.append(f"{fact.get('id')}: failed 원천 사용 시 caveat 누락")
        denominator_fact = fact.get("denominator_fact")
        if denominator_fact and denominator_fact not in fact_ids:
            failures.append(f"{fact.get('id')}: denominator_fact 미해소 {denominator_fact}")
        ratio_like = fact.get("unit") in {"%", "％", "percent"} or bool(
            RATIO_RE.search(fact.get("label", ""))
        )
        if ratio_like and fact.get("denominator") is None and not denominator_fact:
            failures.append(f"{fact.get('id')}: 비율 fact의 분모 누락")
        if ratio_like and not fact.get("comparison_basis"):
            failures.append(f"{fact.get('id')}: 비율 fact의 comparison_basis 누락")

    if failures:
        for detail in failures:
            items.append(_item("source_fact_integrity", "fail", detail, owner="ppt-curator"))
    else:
        items.append(_item("source_fact_integrity", "pass", "source·fact 내부 참조와 비율 제약 일치"))


def _append_fact_value_checks(source: dict[str, Any], items: list[dict[str, Any]]) -> None:
    source_by_id = {entry["id"]: entry for entry in source.get("sources", [])}
    failures: list[str] = []
    unchecked: list[str] = []
    for fact in source.get("facts", []):
        source_entry = source_by_id.get(fact.get("source_id"))
        if source_entry is None:
            continue
        path = _source_asset_path(source_entry["path"])
        try:
            actual = measure_value(path, fact["locator"])
        except (OSError, KeyError, ValueError, json.JSONDecodeError) as error:
            unchecked.append(f"{fact['id']}: {error}")
            continue
        if actual != fact.get("value"):
            failures.append(f"{fact['id']}: 원천 {actual!r}, 매니페스트 {fact.get('value')!r}")

    for detail in failures:
        items.append(_item("fact_value_check", "fail", detail, owner="ppt-curator"))
    for detail in unchecked:
        items.append(_item("fact_value_check", "unchecked", detail))
    if not failures and not unchecked:
        items.append(_item("fact_value_check", "pass", "facts[].value 전건 원천 재측정 일치"))


def _append_zip_and_slide_count(pptx_path: Path, deck: dict[str, Any], items: list[dict[str, Any]]) -> None:
    if not pptx_path.is_file():
        items.append(_item("zip_integrity", "unchecked", "deck.pptx 없음"))
        items.append(_item("slide_count", "unchecked", "deck.pptx 없음"))
        return
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            corrupt = archive.testzip()
    except zipfile.BadZipFile:
        items.append(_item("zip_integrity", "fail", "PPTX ZIP 손상", owner="ppt-renderer"))
        items.append(_item("slide_count", "unchecked", "PPTX를 열 수 없음"))
        return
    if corrupt:
        items.append(_item("zip_integrity", "fail", f"손상 엔트리: {corrupt}", owner="ppt-renderer"))
    else:
        items.append(_item("zip_integrity", "pass", "PPTX ZIP 무결성 통과"))

    try:
        actual = len(Presentation(pptx_path).slides)
    except Exception as error:  # python-pptx 파서가 보고하는 다양한 패키지 오류
        items.append(_item("slide_count", "fail", f"PPTX 파싱 실패: {error}", owner="ppt-renderer"))
        return
    expected = len(deck.get("slides", []))
    if actual == expected:
        items.append(_item("slide_count", "pass", f"슬라이드 수 {actual}장 일치"))
    else:
        items.append(
            _item(
                "slide_count",
                "fail",
                f"실제 {actual}장, 스펙 {expected}장",
                owner="ppt-renderer",
            )
        )


def _normalized_text(value: str) -> str:
    return " ".join(value.split())


def _append_pptx_text_integrity(
    pptx_path: Path,
    deck: dict[str, Any],
    items: list[dict[str, Any]],
) -> None:
    if not pptx_path.is_file():
        items.append(_item("pptx_text_integrity", "unchecked", "deck.pptx 없음"))
        return
    try:
        presentation = Presentation(pptx_path)
    except Exception as error:
        items.append(_item("pptx_text_integrity", "unchecked", f"PPTX 파싱 불가: {error}"))
        return

    failures: list[tuple[int, str]] = []
    for index, slide_spec in enumerate(deck.get("slides", [])):
        if index >= len(presentation.slides):
            break
        actual_text = _normalized_text(
            "\n".join(
                shape.text
                for shape in presentation.slides[index].shapes
                if getattr(shape, "has_text_frame", False)
            )
        )
        expected: list[str] = []
        if not (
            slide_spec.get("role") == "cover"
            and slide_spec.get("layout") == "cover-asym"
        ):
            expected.append(slide_spec.get("title", ""))
        expected.extend(body.get("text", "") for body in slide_spec.get("body", []))
        expected.extend(
            visual.get("caption", "") or "" for visual in slide_spec.get("visuals", [])
        )
        for text in expected:
            normalized = _normalized_text(text)
            if normalized and normalized not in actual_text:
                failures.append((slide_spec["no"], text))

    if failures:
        for slide_no, text in failures:
            items.append(
                _item(
                    "pptx_text_integrity",
                    "fail",
                    f"PPTX에서 스펙 텍스트를 찾지 못함: {text}",
                    slide_no=slide_no,
                    owner="ppt-renderer",
                )
            )
    else:
        items.append(
            _item(
                "pptx_text_integrity",
                "pass",
                "화면 표시 대상 제목·본문 텍스트 전건 PPTX 수록",
            )
        )


def _append_artifact_identity(
    source: dict[str, Any],
    deck: dict[str, Any],
    build_report: dict[str, Any],
    items: list[dict[str, Any]],
) -> None:
    ids = {
        "01_source.json": source.get("deck_id"),
        "02_deck.json": deck.get("deck_id"),
        "_build_report.json": build_report.get("deck_id"),
    }
    failures: list[str] = []
    if len(set(ids.values())) != 1:
        failures.append(", ".join(f"{name}={value}" for name, value in ids.items()))
    if build_report.get("source_sha256") != deck.get("source_sha256"):
        failures.append("_build_report.json의 source_sha256이 02_deck.json과 불일치")
    if failures:
        for detail in failures:
            items.append(_item("artifact_identity", "fail", detail, owner="ppt-renderer"))
    else:
        items.append(_item("artifact_identity", "pass", "deck_id·상류 해시 전달 일치"))


def _append_hash_checks(
    source_path: Path,
    deck_path: Path,
    deck: dict[str, Any],
    build_report: dict[str, Any],
    items: list[dict[str, Any]],
) -> None:
    source_hash = sha256_file(source_path)
    if deck.get("source_sha256") == source_hash:
        items.append(_item("upstream_hash", "pass", "01_source.json 해시 일치"))
    else:
        items.append(_item("upstream_hash", "fail", "01_source.json 해시 불일치", owner="ppt-curator"))

    spec_hash = sha256_file(deck_path)
    if build_report.get("spec_sha256") == spec_hash:
        items.append(_item("spec_hash", "pass", "02_deck.json 해시 일치"))
    else:
        items.append(_item("spec_hash", "fail", "02_deck.json 해시 불일치", owner="ppt-renderer"))


def _append_render_checks(
    run_dir: Path,
    deck: dict[str, Any],
    build_report: dict[str, Any],
    items: list[dict[str, Any]],
) -> None:
    expected = len(deck.get("slides", []))
    build_dir = run_dir / "build"
    expected_previews = [build_dir / "previews" / f"s{number:02d}.png" for number in range(1, expected + 1)]
    reported_previews = [
        Path(path).resolve()
        for path in build_report.get("artifacts", {}).get("previews", [])
    ]
    valid_previews = 0
    for path in expected_previews:
        try:
            with Image.open(path) as image:
                image.verify()
                if image.width > 0 and image.height > 0:
                    valid_previews += 1
        except (OSError, ValueError):
            continue
    expected_resolved = [path.resolve() for path in expected_previews]
    if reported_previews == expected_resolved and valid_previews == expected:
        items.append(_item("previews_present", "pass", f"미리보기 {valid_previews}/{expected}장"))
    else:
        items.append(
            _item(
                "previews_present",
                "fail",
                f"고정 경로의 유효 미리보기 {valid_previews}/{expected}장",
                owner="ppt-renderer",
            )
        )

    artifacts = build_report.get("artifacts", {})
    reported_pptx = Path(artifacts.get("pptx", "")).resolve()
    reported_pdf = artifacts.get("pdf")
    pdf_matches = reported_pdf is None or Path(reported_pdf).resolve() == (build_dir / "deck.pdf").resolve()
    if reported_pptx == (build_dir / "deck.pptx").resolve() and pdf_matches:
        items.append(_item("artifact_paths", "pass", "빌드 산출물 경로가 고정 경계와 일치"))
    else:
        items.append(_item("artifact_paths", "fail", "빌드 리포트 산출물 경로 이탈", owner="ppt-renderer"))

    font = build_report.get("font", {})
    if font.get("fallback"):
        items.append(
            _item(
                "font_fallback",
                "warn",
                f"{font.get('requested')} → {font.get('used')}: {font.get('reason')}",
                owner="ppt-renderer",
            )
        )
    else:
        items.append(_item("font_fallback", "pass", f"요청 폰트 사용: {font.get('used')}"))

    visual = build_report.get("visual_check", {})
    slides_seen = visual.get("slides_seen", [])
    valid_numbers = all(number in range(1, expected + 1) for number in slides_seen)
    if visual.get("done") and slides_seen and valid_numbers and visual.get("note"):
        items.append(_item("visual_check", "pass", f"육안확인 슬라이드 {visual['slides_seen']}"))
    else:
        items.append(_item("visual_check", "fail", "렌더러 육안확인 미기록", owner="ppt-renderer"))


def _append_overflow_check(deck: dict[str, Any], items: list[dict[str, Any]]) -> None:
    candidates: list[tuple[int, str]] = []
    for slide in deck.get("slides", []):
        if len(slide.get("title", "")) > 60:
            candidates.append((slide["no"], "제목 60자 초과"))
        for body in slide.get("body", []):
            if len(body.get("text", "")) > 90:
                candidates.append((slide["no"], "본문 항목 90자 초과"))
    if candidates:
        for slide_no, detail in candidates:
            items.append(_item("text_overflow", "warn", detail, slide_no=slide_no, owner="ppt-writer"))
    else:
        items.append(_item("text_overflow", "pass", "길이 기반 오버플로 후보 없음"))


def _append_closing_check(source: dict[str, Any], deck: dict[str, Any], items: list[dict[str, Any]]) -> None:
    limitations = source.get("limitations", [])
    has_closing = any(slide.get("role") == "closing" for slide in deck.get("slides", []))
    if limitations and not has_closing:
        items.append(
            _item(
                "closing_on_limits",
                "fail",
                "limitations가 있으나 closing 슬라이드 없음",
                owner="ppt-writer",
            )
        )
    else:
        items.append(_item("closing_on_limits", "pass", "한계와 closing 대응"))


def _write_markdown(path: Path, report: dict[str, Any]) -> None:
    lines = [
        f"# PPT QA — {report['deck_id']}",
        "",
        f"- 판정: **{report['verdict']}**",
        (
            f"- 통과 {report['summary']['passed']} · 실패 {report['summary']['failed']} · "
            f"경고 {report['summary']['warned']} · 미검사 {report['summary']['unchecked']}"
        ),
        "",
        "## 검사 항목",
        "",
    ]
    for item in report["items"]:
        location = f" · s{item['slide_no']:02d}" if item.get("slide_no") else ""
        owner = f" · 반려 {item['owner']}" if item.get("owner") else ""
        detail = f" — {item['detail']}" if item.get("detail") else ""
        lines.append(f"- `{item['status']}` {item['name']}{location}{owner}{detail}")
    lines.extend(["", "## 사각지대", ""])
    lines.extend(f"- {value}" for value in report["blind_spots"])
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def verify(
    run_dir: Path | str,
    as_of: str | None = None,
    *,
    adversarial_review_note: str | None = None,
) -> dict[str, Any]:
    """run 디렉터리를 읽기 전용 검사하고 QA 리포트 두 파일만 쓴다."""
    run = Path(run_dir).resolve()
    source_path = run / "01_source.json"
    deck_path = run / "02_deck.json"
    build_path = run / "build" / "_build_report.json"
    pptx_path = run / "build" / "deck.pptx"

    items: list[dict[str, Any]] = []

    source, source_ok = _load_and_validate_json(
        source_path,
        SOURCE_SCHEMA,
        "source_schema",
        "ppt-curator",
        items,
    )
    deck, deck_ok = _load_and_validate_json(
        deck_path,
        DECK_SCHEMA,
        "deck_schema",
        "ppt-writer",
        items,
    )
    build_report, build_ok = _load_and_validate_json(
        build_path,
        BUILD_SCHEMA,
        "build_schema",
        "ppt-renderer",
        items,
    )

    if source_ok and source is not None:
        _append_source_asset_hashes(source, items)
        _append_source_fact_integrity(source, items)
        _append_fact_value_checks(source, items)
    else:
        _append_unchecked(items, SOURCE_CHECKS, "01_source.json 구조 미충족")

    if deck is not None:
        _append_title_checks(deck, items)
    else:
        items.append(_item("title_style", "unchecked", "02_deck.json 구조 미충족"))

    if deck_ok and deck is not None:
        _append_slide_structure(deck, items)
        _append_bullet_hierarchy(deck, items)
        threshold = float(_load_json(QA_CONTRACT)["기본값"]["body_nominal_ratio_min"])
        _append_body_checks(deck, items, threshold)
        _append_overflow_check(deck, items)
    else:
        _append_unchecked(items, DECK_CHECKS[1:], "02_deck.json 구조 미충족")

    if source_ok and deck_ok and source is not None and deck is not None:
        _append_fact_checks(source, deck, items)
        _append_closing_check(source, deck, items)
    else:
        _append_unchecked(items, FACT_CHECKS, "source 또는 deck 구조 미충족")

    if deck_ok and deck is not None:
        _append_zip_and_slide_count(pptx_path, deck, items)
        _append_pptx_text_integrity(pptx_path, deck, items)
    else:
        _append_unchecked(items, PPTX_CHECKS, "02_deck.json 구조 미충족")

    if source_ok and deck_ok and build_ok and source is not None and deck is not None and build_report is not None:
        _append_artifact_identity(source, deck, build_report, items)
        _append_hash_checks(source_path, deck_path, deck, build_report, items)
        _append_render_checks(run, deck, build_report, items)
    else:
        _append_unchecked(items, BUILD_CHECKS, "source·deck·build 구조 중 하나가 미충족")

    current_verified = {
        "pptx_sha256": sha256_file(pptx_path) if pptx_path.is_file() else None,
        "spec_sha256": sha256_file(deck_path) if deck_path.is_file() else None,
        "source_sha256": sha256_file(source_path) if source_path.is_file() else None,
    }
    if adversarial_review_note and adversarial_review_note.strip():
        items.append(_item("adversarial_review", "pass", adversarial_review_note.strip()))
    else:
        items.append(_item("adversarial_review", "unchecked", "원천 반박검증 기록 없음"))

    counts = Counter(item["status"] for item in items)
    summary = {
        "passed": counts["pass"],
        "failed": counts["fail"],
        "warned": counts["warn"],
        "unchecked": counts["unchecked"],
    }
    report = {
        "deck_id": (deck or source or {}).get("deck_id", run.name),
        "as_of": as_of or date.today().isoformat(),
        "verified": {
            **current_verified,
        },
        "summary": summary,
        "items": items,
        "verdict": "accept" if summary["failed"] == 0 and summary["unchecked"] == 0 else "reject",
        "blind_spots": BLIND_SPOTS,
        "notes": [],
    }
    jsonschema.validate(report, _load_json(QA_SCHEMA))
    run.mkdir(parents=True, exist_ok=True)
    (run / "03_qa.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    _write_markdown(run / "03_qa.md", report)
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("run_dir", nargs="?", type=Path, help="output/ppt/<deck-id> 경로. 생략 시 전 덱 검사")
    parser.add_argument("--as-of", help="QA 기준일 YYYY-MM-DD")
    parser.add_argument("--adversarial-review-note", help="원천 반박검증의 대상·범위·결과")
    args = parser.parse_args()

    run_dirs = [args.run_dir] if args.run_dir else sorted(
        path.parent.parent
        for path in (ROOT / "output/ppt").glob("*/build/_build_report.json")
    )
    if not run_dirs:
        print("[2] 검증할 output/ppt/<deck-id> 빌드 없음")
        return 2

    output_root = (ROOT / "output/ppt").resolve()
    if any(not run_dir.resolve().is_relative_to(output_root) for run_dir in run_dirs):
        parser.error("run_dir은 output/ppt/<deck-id> 아래여야 함")

    reports = [
        verify(
            run_dir,
            args.as_of,
            adversarial_review_note=(
                args.adversarial_review_note or os.environ.get("PPT_ADVERSARIAL_REVIEW_NOTE")
            ),
        )
        for run_dir in run_dirs
    ]
    print(json.dumps(reports[0] if len(reports) == 1 else reports, ensure_ascii=False, indent=2))
    return 0 if all(report["verdict"] == "accept" for report in reports) else 1


if __name__ == "__main__":
    raise SystemExit(main())
