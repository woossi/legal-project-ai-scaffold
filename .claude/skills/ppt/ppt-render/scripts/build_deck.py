#!/usr/bin/env python3
"""02_deck.json을 로컬 python-pptx 덱과 빌드 자기보고로 변환한다."""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Iterable

import jsonschema
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SCRIPT = Path(__file__).resolve()
ROOT = next(parent for parent in SCRIPT.parents if (parent / "AGENTS.md").is_file())
DECK_SCHEMA = ROOT / ".claude/skills/ppt/ppt-deck/contract/deck.schema.json"
BUILD_SCHEMA = ROOT / ".claude/skills/ppt/ppt-render/contract/build.schema.json"
DESIGN_TOKENS = ROOT / "docs/ppt/design-tokens.json"

DEFAULT_STYLE = {
    "canvas_width": 13.333,
    "canvas_height": 7.5,
    "primary": "4E83F9",
    "primary_strong": "1F4FB0",
    "primary_soft": "EAF0FD",
    "primary_pale": "F6F9FF",
    "surface": "FFFFFF",
    "surface_subtle": "F6F8F9",
    "text": "242424",
    "muted": "6B7280",
    "border": "D9D9D9",
    "feedback": "E17445",
    "feedback_soft": "FFF4ED",
    "font": "Apple SD Gothic Neo",
    "cover_title_pt": 54,
    "title_pt": 30,
    "body_pt": 20,
    "body_small_pt": 16,
    "eyebrow_pt": 15,
    "caption_pt": 13,
}


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _hex(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#"))


def _token_value(tokens: dict[str, Any], *keys: str) -> Any:
    current: Any = tokens
    for key in keys:
        current = current[key]
    return current.get("$value") if isinstance(current, dict) and "$value" in current else current


def _token_hex(tokens: dict[str, Any], *keys: str) -> str:
    return str(_token_value(tokens, *keys)).lstrip("#")


def _token_pt(tokens: dict[str, Any], *keys: str) -> float:
    return float(str(_token_value(tokens, *keys)).removesuffix("pt"))


def load_style(profile: str) -> dict[str, Any]:
    """프로필을 렌더러 내부의 작은 스타일 표로 바꾼다."""
    style = dict(DEFAULT_STYLE)
    if profile == "none":
        return style
    if profile != "docs-ppt-20260813":
        raise ValueError(f"지원하지 않는 profile: {profile}")
    if not DESIGN_TOKENS.is_file():
        raise FileNotFoundError(DESIGN_TOKENS)

    tokens = json.loads(DESIGN_TOKENS.read_text(encoding="utf-8"))
    style.update(
        {
            "canvas_width": 20.0,
            "canvas_height": 11.25,
            "primary": _token_hex(tokens, "color", "primary"),
            "primary_strong": _token_hex(tokens, "color", "primary-strong"),
            "primary_soft": _token_hex(tokens, "color", "primary-soft"),
            "primary_pale": _token_hex(tokens, "color", "primary-pale"),
            "surface": _token_hex(tokens, "color", "surface"),
            "surface_subtle": _token_hex(tokens, "color", "surface-subtle"),
            "text": _token_hex(tokens, "color", "text", "primary"),
            "muted": _token_hex(tokens, "color", "text", "muted"),
            "border": _token_hex(tokens, "color", "border"),
            "feedback": _token_hex(tokens, "color", "feedback"),
            "feedback_soft": _token_hex(tokens, "color", "feedback-soft"),
            "font": _token_value(tokens, "typography", "font-family", "content-medium"),
            "cover_title_pt": _token_pt(tokens, "typography", "font-size", "cover-title"),
            "title_pt": _token_pt(tokens, "typography", "font-size", "slide-title"),
            "body_pt": _token_pt(tokens, "typography", "font-size", "body"),
            "body_small_pt": _token_pt(tokens, "typography", "font-size", "body-small"),
            "eyebrow_pt": _token_pt(tokens, "typography", "font-size", "eyebrow"),
            "caption_pt": _token_pt(tokens, "typography", "font-size", "caption"),
        }
    )
    return style


def apply_deck_overrides(style: dict[str, Any], spec: dict[str, Any]) -> dict[str, Any]:
    """프로필은 유지하고, 사용자 지시가 있는 덱 단위 글꼴만 덮어쓴다."""
    resolved = dict(style)
    if spec.get("font_family"):
        resolved["font"] = str(spec["font_family"])
    return resolved


def _font_file_patterns(family: str) -> tuple[str, ...]:
    normalized = family.lower().replace(" ", "")
    if "kopub" in normalized:
        return ("*KoPub*", "*KOPUB*")
    if "nanum" in normalized:
        return ("*Nanum*",)
    if "applesdgothic" in normalized:
        return ("*AppleSDGothic*",)
    return (f"*{family.replace(' ', '*')}*",)


def font_installed(family: str) -> bool:
    roots = (
        Path("/System/Library/Fonts"),
        Path("/Library/Fonts"),
        Path.home() / "Library/Fonts",
    )
    for root in roots:
        if not root.is_dir():
            continue
        for pattern in _font_file_patterns(family):
            if next(root.rglob(pattern), None) is not None:
                return True
    return False


def select_font(requested: str) -> dict[str, Any]:
    if font_installed(requested):
        return {"requested": requested, "used": requested, "fallback": False, "reason": None}
    for candidate in ("Apple SD Gothic Neo", "NanumGothic"):
        if font_installed(candidate):
            return {
                "requested": requested,
                "used": candidate,
                "fallback": True,
                "reason": f"요청 폰트 미설치: {requested}",
            }
    return {
        "requested": requested,
        "used": requested,
        "fallback": True,
        "reason": "확인 가능한 한글 폴백 폰트 없음",
    }


def _reset_text_frame(
    frame: Any,
    *,
    vertical_anchor: MSO_ANCHOR | None = None,
) -> None:
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    if vertical_anchor is not None:
        frame.vertical_anchor = vertical_anchor


def _body_item_text(item: dict[str, Any]) -> str:
    level = int(item.get("level", 1))
    kind = item.get("kind")
    if kind:
        return f"• {kind} — {item['text']}"
    if level >= 2:
        return f"◦ {item['text']}"
    return item["text"]


def _add_textbox(
    slide: Any,
    text: str,
    box: tuple[float, float, float, float],
    *,
    font: str,
    size: float,
    color: str,
    bold: bool = False,
    underline: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    shape = slide.shapes.add_textbox(*(Inches(value) for value in box))
    frame = shape.text_frame
    _reset_text_frame(frame, vertical_anchor=MSO_ANCHOR.MIDDLE)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.underline = underline
    run.font.color.rgb = _hex(color)
    return shape


def _add_label_content_textbox(
    slide: Any,
    text: str,
    box: tuple[float, float, float, float],
    *,
    font: str,
    label_size: float,
    content_size: float,
    label_color: str,
    content_color: str,
) -> Any:
    """`표지 — 설명`을 한 도형에 유지하면서 표지만 강조한다."""
    label, content = text.split(" — ", 1) if " — " in text else (text, "")
    shape = slide.shapes.add_textbox(*(Inches(value) for value in box))
    frame = shape.text_frame
    _reset_text_frame(frame, vertical_anchor=MSO_ANCHOR.MIDDLE)
    paragraph = frame.paragraphs[0]
    label_run = paragraph.add_run()
    label_run.text = f"{label} — " if content else label
    label_run.font.name = font
    label_run.font.size = Pt(label_size)
    label_run.font.bold = True
    label_run.font.color.rgb = _hex(label_color)
    if content:
        content_run = paragraph.add_run()
        content_run.text = content
        content_run.font.name = font
        content_run.font.size = Pt(content_size)
        content_run.font.color.rgb = _hex(content_color)
    return shape


def _add_shape(
    slide: Any,
    shape_type: MSO_SHAPE,
    box: tuple[float, float, float, float],
    *,
    fill: str,
    line: str | None = None,
    rotation: float = 0,
) -> Any:
    """스타일 색만 쓰는 편집 가능한 PowerPoint 도형을 추가한다."""
    shape = slide.shapes.add_shape(
        shape_type,
        *(Inches(value) for value in box),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(fill)
    if line:
        shape.line.color.rgb = _hex(line)
    else:
        shape.line.fill.background()
    shape.rotation = rotation
    return shape


def _add_rule(
    slide: Any,
    box: tuple[float, float, float, float],
    color: str,
    *,
    rotation: float = 0,
) -> Any:
    return _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        box,
        fill=color,
        rotation=rotation,
    )


def _add_badge_icon(
    slide: Any,
    icon: str,
    box: tuple[float, float, float, float],
    style: dict[str, Any],
) -> Any:
    """check·arrow·alert를 글꼴 아이콘이 아닌 기본 도형으로 그린다."""
    x, y, width, height = box
    feedback = icon == "alert"
    ink = style["feedback"] if feedback else style["primary"]
    soft = style["feedback_soft"] if feedback else style["primary_soft"]
    badge = _add_shape(
        slide,
        MSO_SHAPE.OVAL,
        box,
        fill=soft,
        line=ink,
    )
    if icon == "check":
        _add_rule(
            slide,
            (x + width * 0.25, y + height * 0.52, width * 0.24, height * 0.10),
            ink,
            rotation=42,
        )
        _add_rule(
            slide,
            (x + width * 0.39, y + height * 0.41, width * 0.39, height * 0.10),
            ink,
            rotation=-43,
        )
    elif icon == "arrow":
        _add_shape(
            slide,
            MSO_SHAPE.RIGHT_ARROW,
            (x + width * 0.20, y + height * 0.34, width * 0.60, height * 0.31),
            fill=ink,
        )
    elif icon == "alert":
        _add_rule(
            slide,
            (x + width * 0.46, y + height * 0.22, width * 0.09, height * 0.38),
            ink,
        )
        _add_shape(
            slide,
            MSO_SHAPE.OVAL,
            (x + width * 0.44, y + height * 0.69, width * 0.13, height * 0.13),
            fill=ink,
        )
    return badge


def _add_role_badge(
    slide: Any,
    role: str,
    box: tuple[float, float, float, float],
    style: dict[str, Any],
) -> None:
    x, y, width, height = box
    _add_shape(
        slide,
        MSO_SHAPE.OVAL,
        box,
        fill=style["primary_soft"],
        line=style["primary"],
    )
    silhouettes = {
        "task": MSO_SHAPE.FLOWCHART_DOCUMENT,
        "bridge": MSO_SHAPE.RIGHT_ARROW,
        "appendix": MSO_SHAPE.HEXAGON,
        "closing": MSO_SHAPE.OVAL,
    }
    _add_shape(
        slide,
        silhouettes.get(role, MSO_SHAPE.FLOWCHART_PROCESS),
        (x + width * 0.25, y + height * 0.28, width * 0.50, height * 0.44),
        fill=style["primary"],
    )


def _add_cover_chrome(
    slide: Any,
    *,
    as_of: str,
    style: dict[str, Any],
    font: str,
) -> None:
    """파란 표지 위에 날짜·범위 스트립과 절제된 연결 모티프를 둔다."""
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    unit = max(0.65, min(1.0, canvas_h / 11.25))
    _add_textbox(
        slide,
        f"기준일 {as_of}  ·  프로젝트 진행 결과",
        (0.9, 0.48, canvas_w - 1.8, 0.35),
        font=font,
        size=style["caption_pt"],
        color=style["surface"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    motif = (
        (canvas_w - 3.00 * unit, 1.15 * unit, 0.54 * unit),
        (canvas_w - 1.75 * unit, 2.00 * unit, 0.34 * unit),
        (canvas_w - 3.70 * unit, 3.10 * unit, 0.26 * unit),
        (0.95 * unit, canvas_h - 2.20 * unit, 0.42 * unit),
        (2.10 * unit, canvas_h - 1.45 * unit, 0.25 * unit),
    )
    for x, y, diameter in motif:
        _add_shape(
            slide,
            MSO_SHAPE.OVAL,
            (x, y, diameter, diameter),
            fill=style["primary"],
            line=style["primary_soft"],
        )
    _add_rule(
        slide,
        (canvas_w - 2.67 * unit, 1.68 * unit, 1.15 * unit, 0.035 * unit),
        style["primary_soft"],
        rotation=28,
    )
    _add_rule(
        slide,
        (1.23 * unit, canvas_h - 1.88 * unit, 1.02 * unit, 0.035 * unit),
        style["primary_soft"],
        rotation=30,
    )


def _add_noncover_chrome(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
) -> tuple[float, float]:
    """역할 표식·제목 배지·accent rule을 공통 상단 chrome으로 렌더한다."""
    role_labels = {
        "task": "작업 결과",
        "bridge": "흐름 연결",
        "appendix": "근거 자료",
        "closing": "마무리",
    }
    section_labels = {
        "bridge": "전체 흐름",
        "appendix": "근거 부록",
        "closing": "결과 정리",
    }
    role = slide_spec["role"]
    section = slide_spec.get("work_boundary") or section_labels.get(role, "프로젝트 결과")
    section = {
        "작업 1 — 수집 단계(자료 수집 및 정제)": "자료 수집·정제",
        "작업 2 — 법 적용 확인": "법 적용 확인",
        "작업 3 — 스키마 설계 및 구현": "온톨로지 설계·구현",
    }.get(section, section)
    _add_textbox(
        slide,
        f"{role_labels.get(role, role)}  ·  {section}",
        (0.92, 0.16, style["canvas_width"] - 1.84, 0.27),
        font=font,
        size=style["eyebrow_pt"],
        color=style["muted"],
        bold=True,
    )
    if slide_spec.get("layout") != "glossary":
        _add_textbox(
            slide,
            "용어 설명 ↗",
            (style["canvas_width"] - 2.28, 0.15, 1.36, 0.29),
            font=font,
            size=style["caption_pt"],
            color=style["primary"],
            bold=True,
            underline=True,
            align=PP_ALIGN.RIGHT,
        )
    _add_role_badge(slide, role, (0.92, 0.64, 0.46, 0.46), style)
    _add_rule(
        slide,
        (0.92, 1.47, style["canvas_width"] - 1.84, 0.025),
        style["primary"],
    )
    return 1.52, style["canvas_width"] - 2.44


def _add_body(
    slide: Any,
    items: list[dict[str, Any]],
    box: tuple[float, float, float, float],
    style: dict[str, Any],
    font: str,
    *,
    color: str | None = None,
) -> None:
    shape = slide.shapes.add_textbox(*(Inches(value) for value in box))
    frame = shape.text_frame
    _reset_text_frame(frame)

    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        level = int(item.get("level", 1))
        paragraph.text = _body_item_text(item)
        paragraph.level = max(0, level - 1)
        paragraph.space_after = Pt(10 if level == 1 else 6)
        paragraph.font.name = font
        paragraph.font.size = Pt(style["body_pt"])
        paragraph.font.bold = bool(item.get("emphasis", False))
        paragraph.font.color.rgb = _hex(color or style["text"])


def _add_stage_rail(
    slide: Any,
    stage: int,
    style: dict[str, Any],
    font: str,
) -> None:
    labels = ("수집", "내용 보존", "적용 범위", "근거 연결", "누락 검증", "전국 확대")
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    left = 1.25
    right = canvas_w - 1.25
    rail_y = canvas_h - 0.95
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(rail_y),
        Inches(right - left),
        Inches(0.025),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _hex(style["border"])
    line.line.fill.background()

    step = (right - left) / (len(labels) - 1)
    for index, label in enumerate(labels, start=1):
        center_x = left + (index - 1) * step
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(center_x - 0.09),
            Inches(rail_y - 0.08),
            Inches(0.18),
            Inches(0.18),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _hex(
            style["primary"] if index == stage else style["surface"]
        )
        dot.line.color.rgb = _hex(
            style["primary"] if index == stage else style["border"]
        )
        _add_textbox(
            slide,
            label,
            (center_x - 0.65, rail_y + 0.18, 1.3, 0.32),
            font=font,
            size=style["caption_pt"],
            color=style["primary"] if index == stage else style["muted"],
            bold=index == stage,
            align=PP_ALIGN.CENTER,
        )


def _resolve_asset(path_value: str, spec_path: Path) -> Path:
    candidate = Path(path_value)
    if candidate.is_absolute():
        return candidate
    for base in (ROOT, spec_path.parent):
        resolved = base / candidate
        if resolved.is_file():
            return resolved
    raise FileNotFoundError(path_value)


def _add_image(
    slide: Any,
    asset: Path,
    box: tuple[float, float, float, float],
    *,
    crop_bottom: float = 0.0,
) -> Any:
    x, y, width, height = box
    with Image.open(asset) as image:
        ratio = image.width / image.height
    box_ratio = width / height
    if ratio >= box_ratio:
        fitted_width = width
        fitted_height = width / ratio
        fitted_x = x
        fitted_y = y + (height - fitted_height) / 2
    else:
        fitted_height = height
        fitted_width = height * ratio
        fitted_x = x + (width - fitted_width) / 2
        fitted_y = y
    picture = slide.shapes.add_picture(
        str(asset),
        Inches(fitted_x),
        Inches(fitted_y),
        Inches(fitted_width),
        Inches(fitted_height),
    )
    if crop_bottom:
        picture.crop_bottom = crop_bottom
    return picture


def _add_table(
    slide: Any,
    asset: Path,
    box: tuple[float, float, float, float],
    style: dict[str, Any],
    font: str,
) -> None:
    with asset.open(encoding="utf-8-sig", newline="") as stream:
        rows = list(csv.reader(stream))
    if not rows:
        return
    column_count = max(len(row) for row in rows)
    if len(rows) > 11 or column_count > 6:
        raise ValueError(f"표 자산 한도 초과: {len(rows)}행 × {column_count}열, 허용 11행 × 6열")
    table_shape = slide.shapes.add_table(
        len(rows),
        column_count,
        *(Inches(value) for value in box),
    )
    table = table_shape.table
    for row_index, row in enumerate(rows):
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            cell.text = row[column_index] if column_index < len(row) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = font
                paragraph.font.size = Pt(style["caption_pt"])
                paragraph.font.bold = row_index == 0


def _group_decision_body(
    items: list[dict[str, Any]],
) -> dict[str, list[dict[str, Any]]]:
    groups = {kind: [] for kind in ("결과", "결론", "남은 한계")}
    current: str | None = None
    for item in items:
        kind = item.get("kind")
        if int(item.get("level", 1)) == 1 and kind in groups:
            current = str(kind)
        if current:
            groups[current].append(item)
    return groups


def _add_column_copy(
    slide: Any,
    items: list[dict[str, Any]],
    box: tuple[float, float, float, float],
    style: dict[str, Any],
    font: str,
    *,
    dense: bool = False,
    accent: str | None = None,
) -> None:
    shape = slide.shapes.add_textbox(*(Inches(value) for value in box))
    frame = shape.text_frame
    _reset_text_frame(frame)
    minimum = max(16.0, float(style["body_small_pt"]))
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        child = int(item.get("level", 1)) >= 2
        paragraph.text = f"• {item['text']}" if child else item["text"]
        paragraph.space_after = Pt(5 if dense and child else 8 if child else 10)
        paragraph.font.name = font
        paragraph.font.size = Pt(
            minimum if child or dense else max(minimum + 2, style["body_pt"])
        )
        paragraph.font.bold = not child or bool(item.get("emphasis", False))
        paragraph.font.underline = bool(item.get("emphasis", False))
        paragraph.font.color.rgb = _hex(style["text"])
        if not child and accent:
            paragraph.font.color.rgb = _hex(accent)


def _friendly_kind_label(kind: str) -> str:
    return {
        "결과": "무엇을 했나",
        "결론": "왜 했나",
        "남은 한계": "다음 확인",
    }.get(kind, kind)


def _decision_icon(kind: str) -> str:
    return {
        "결과": "check",
        "결론": "arrow",
        "남은 한계": "alert",
    }.get(kind, "arrow")


def _decision_accent(kind: str, style: dict[str, Any]) -> str:
    return {
        "결과": style["primary"],
        "결론": style["primary_strong"],
        "남은 한계": style["feedback"],
    }.get(kind, style["primary"])


def _decision_fill(kind: str, style: dict[str, Any]) -> str:
    return style["feedback_soft"] if kind == "남은 한계" else style["primary_pale"]


def _add_decision_block(
    slide: Any,
    kind: str,
    items: list[dict[str, Any]],
    box: tuple[float, float, float, float],
    style: dict[str, Any],
    font: str,
    *,
    dense: bool = False,
    mirror: bool = False,
) -> None:
    """같은 의미 체계를 유지하면서 블록의 무게중심만 바꾼다."""
    x, y, width, height = box
    accent = _decision_accent(kind, style)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, box, fill=_decision_fill(kind, style))
    if mirror:
        _add_rule(slide, (x + width - 0.055, y, 0.055, height), accent)
    else:
        _add_rule(slide, (x, y, 0.055, height), accent)
    _add_rule(slide, (x + 0.25, y + height - 0.20, width - 0.50, 0.035), style["border"])
    watermark_shape = {
        "결과": MSO_SHAPE.FLOWCHART_DOCUMENT,
        "결론": MSO_SHAPE.RIGHT_ARROW,
        "남은 한계": MSO_SHAPE.FLOWCHART_DECISION,
    }.get(kind, MSO_SHAPE.FLOWCHART_PROCESS)
    watermark_size = min(1.45, max(0.85, height * 0.30))
    watermark_x = x + 0.34 if mirror else x + width - watermark_size - 0.34
    _add_shape(
        slide,
        watermark_shape,
        (watermark_x, y + height - watermark_size - 0.38, watermark_size, watermark_size),
        fill=style["surface"],
        line=style["primary_soft"] if kind != "남은 한계" else style["feedback_soft"],
    )

    icon_x = x + width - 0.74 if mirror else x + 0.24
    label_x = x + 0.24 if mirror else x + 0.88
    _add_badge_icon(slide, _decision_icon(kind), (icon_x, y + 0.18, 0.50, 0.50), style)
    _add_textbox(
        slide,
        _friendly_kind_label(kind),
        (label_x, y + 0.16, width - 1.18, 0.36),
        font=font,
        size=style["eyebrow_pt"],
        color=accent,
        bold=True,
        align=PP_ALIGN.RIGHT if mirror else PP_ALIGN.LEFT,
    )
    _add_textbox(
        slide,
        kind,
        (label_x, y + 0.52, width - 1.18, 0.28),
        font=font,
        size=style["caption_pt"],
        color=style["muted"],
        bold=True,
        align=PP_ALIGN.RIGHT if mirror else PP_ALIGN.LEFT,
    )

    if any(item.get("emphasis") for item in items):
        _add_rule(slide, (x + 0.24, y + 0.92, width - 0.48, 0.09), style["feedback_soft"])
    _add_column_copy(
        slide,
        items,
        (x + 0.26, y + 0.90, width - 0.52, height - 1.18),
        style,
        font,
        dense=dense,
        accent=accent,
    )


def _decision_block_boxes(
    slide_spec: dict[str, Any],
    style: dict[str, Any],
) -> dict[str, tuple[float, float, float, float]]:
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage_reserved = 1.45 if slide_spec.get("stage") else 0.85
    left = 0.92
    top = 1.72
    bottom = canvas_h - stage_reserved
    width = canvas_w - left * 2
    height = bottom - top
    gap = 0.30
    variant = int(slide_spec.get("no", 1)) % 3

    if variant == 0:
        result_h = 2.35
        lower_y = top + result_h + gap
        lower_h = height - result_h - gap
        half_w = (width - gap) / 2
        return {
            "결과": (left, top, width, result_h),
            "결론": (left, lower_y, half_w, lower_h),
            "남은 한계": (left + half_w + gap, lower_y, half_w, lower_h),
        }

    major_w = width * 0.44
    minor_w = width - major_w - gap
    upper_h = height * 0.47
    lower_y = top + upper_h + gap
    lower_h = height - upper_h - gap
    if variant == 1:
        major_x = left
        minor_x = left + major_w + gap
    else:
        minor_x = left
        major_x = left + minor_w + gap
    return {
        "결과": (major_x, top, major_w, upper_h),
        "결론": (minor_x, top, minor_w, upper_h),
        "남은 한계": (left, lower_y, width, lower_h),
    }


def _render_decision_columns(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
) -> None:
    """결과·결론·남은 한계를 비대칭 편집 블록으로 나눈다."""
    groups = _group_decision_body(slide_spec.get("body", []))
    boxes = _decision_block_boxes(slide_spec, style)
    for kind in ("결과", "결론", "남은 한계"):
        _add_decision_block(
            slide,
            kind,
            groups[kind],
            boxes[kind],
            style,
            font,
            dense=kind != "결과",
            mirror=boxes[kind][0] > style["canvas_width"] / 2,
        )


def _render_cover_asym(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
    *,
    as_of: str,
) -> None:
    """표지는 보고 기간과 작업 내역만 비대칭으로 배치한다."""
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    period = str(slide_spec.get("period") or as_of).replace("-", ".")
    _add_textbox(
        slide,
        period,
        (0.98, 0.74, canvas_w * 0.45, 0.72),
        font=font,
        size=max(24.0, float(style["title_pt"])),
        color=style["surface"],
        bold=True,
    )
    _add_rule(slide, (0.98, 1.56, canvas_w * 0.30, 0.055), style["primary_soft"])

    body = list(slide_spec.get("body", []))
    positions = (
        (0.98, 2.32, canvas_w * 0.62),
        (canvas_w * 0.23, 5.00, canvas_w * 0.67),
        (canvas_w * 0.08, 7.74, canvas_w * 0.60),
    )
    for index, item in enumerate(body[:3]):
        x, y, width = positions[index]
        text = str(item["text"])
        _add_rule(slide, (x, y, 0.10, 1.54), style["primary_soft"])
        _add_label_content_textbox(
            slide,
            text,
            (x + 0.30, y - 0.02, width - 0.30, 1.48),
            font=font,
            label_size=max(22.0, float(style["body_pt"])),
            content_size=max(16.0, float(style["body_small_pt"])),
            label_color=style["surface"],
            content_color=style["primary_soft"],
        )


def _story_box_map(
    layout: str,
    style: dict[str, Any],
    slide_spec: dict[str, Any],
) -> dict[str, tuple[float, float, float, float]]:
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage_reserved = 1.45 if slide_spec.get("stage") else 0.85
    left = 0.92
    top = 1.72
    bottom = canvas_h - stage_reserved
    content_w = canvas_w - left * 2
    content_h = bottom - top
    gap = 0.28
    main_w = content_w * 0.62
    side_w = content_w - main_w - gap
    if layout == "story-right":
        side_w = content_w * 0.44
        main_w = content_w - side_w - gap
        main_x = left + side_w + gap
        side_x = left
    else:
        main_x = left
        side_x = left + main_w + gap
    return {
        "main": (main_x, top, main_w, content_h),
        "why": (side_x, top, side_w, content_h * 0.32),
        "check": (side_x, top + content_h * 0.34, side_w, content_h * 0.31),
        "next": (side_x, top + content_h * 0.67, side_w, content_h * 0.33),
    }


def _add_labeled_panel(
    slide: Any,
    label: str,
    items: list[dict[str, Any]],
    box: tuple[float, float, float, float],
    style: dict[str, Any],
    font: str,
    *,
    accent: str,
    fill: str,
    dense: bool = True,
    mirror: bool = False,
    spread: bool = False,
    watermark: bool = True,
) -> None:
    x, y, width, height = box
    _add_shape(slide, MSO_SHAPE.RECTANGLE, box, fill=fill, line=style["border"])
    if watermark and not spread:
        watermark_shape = MSO_SHAPE.FLOWCHART_DOCUMENT if accent != style["feedback"] else MSO_SHAPE.FLOWCHART_DECISION
        mark_size = min(1.35, max(0.72, height * 0.24))
        _add_shape(
            slide,
            watermark_shape,
            (x + width - mark_size - 0.22, y + height - mark_size - 0.22, mark_size, mark_size),
            fill=fill,
            line=style["primary_soft"] if accent != style["feedback"] else style["feedback_soft"],
        )
    _add_rule(slide, (x + 0.20, y + 0.56, width - 0.40, 0.04), accent)
    _add_textbox(
        slide,
        label,
        (x + 0.22, y + 0.12, width - 0.44, 0.38),
        font=font,
        size=style["eyebrow_pt"],
        color=accent,
        bold=True,
        align=PP_ALIGN.RIGHT if mirror else PP_ALIGN.LEFT,
    )
    copy_box = (x + 0.22, y + 0.70, width - 0.44, max(0.6, height - 0.88))
    if spread:
        _add_spread_copy(slide, items, copy_box, style, font, accent=accent)
    else:
        _add_column_copy(
            slide,
            items,
            copy_box,
            style,
            font,
            dense=dense,
            accent=accent,
        )


def _add_spread_copy(
    slide: Any,
    items: list[dict[str, Any]],
    box: tuple[float, float, float, float],
    style: dict[str, Any],
    font: str,
    *,
    accent: str,
) -> None:
    """큰 패널의 문장을 같은 높이의 정보 밴드로 펼쳐 빈 공간을 줄인다."""
    if not items:
        return
    x, y, width, height = box
    gap = 0.10
    band_h = max(0.54, (height - gap * (len(items) - 1)) / len(items))
    for index, item in enumerate(items):
        band_y = y + index * (band_h + gap)
        emphasized = bool(item.get("emphasis", False)) or int(item.get("level", 1)) == 1
        fill = style["primary_soft"] if emphasized and accent != style["feedback"] else style["surface"]
        if accent == style["feedback"]:
            fill = style["feedback_soft"] if emphasized else style["surface"]
        _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            (x, band_y, width, band_h),
            fill=fill,
            line=style["border"],
        )
        _add_textbox(
            slide,
            f"{index + 1:02d}",
            (x + 0.14, band_y + 0.10, 0.42, max(0.28, band_h - 0.20)),
            font=font,
            size=style["caption_pt"],
            color=accent,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide,
            str(item["text"]),
            (x + 0.72, band_y + 0.08, width - 0.90, max(0.34, band_h - 0.16)),
            font=font,
            size=max(15.0, float(style["body_small_pt"])),
            color=accent if emphasized else style["text"],
            bold=emphasized,
            underline=bool(item.get("emphasis", False)),
        )


def _render_story_asym(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
    *,
    layout: str,
) -> None:
    groups = _group_decision_body(slide_spec.get("body", []))
    boxes = _story_box_map(layout, style, slide_spec)
    mirror = layout == "story-right"
    _add_labeled_panel(
        slide,
        "무엇을 했나",
        groups["결과"],
        boxes["main"],
        style,
        font,
        accent=style["primary"],
        fill=style["primary_pale"],
        dense=False,
        mirror=mirror,
        spread=True,
    )
    _add_labeled_panel(
        slide,
        "왜 했나",
        groups["결론"][:1],
        boxes["why"],
        style,
        font,
        accent=style["primary_strong"],
        fill=style["surface"],
        mirror=not mirror,
    )
    _add_labeled_panel(
        slide,
        "작업 내용·확인",
        groups["결론"][1:] or groups["결론"],
        boxes["check"],
        style,
        font,
        accent=style["primary"],
        fill=style["surface_subtle"],
        mirror=not mirror,
    )
    _add_labeled_panel(
        slide,
        "다음 방향",
        groups["남은 한계"],
        boxes["next"],
        style,
        font,
        accent=style["feedback"],
        fill=style["feedback_soft"],
        mirror=not mirror,
    )


def _first_number(text: str) -> str | None:
    token = ""
    for char in text:
        if char.isdigit() or (token and char == ","):
            token += char
            continue
        if token and char != ",":
            return token
    return token or None


def _metric_tokens(items: list[dict[str, Any]], limit: int = 3) -> list[str]:
    """문장에 이미 결박된 수치와 단위를 숫자 칩에 재사용한다."""
    pattern = re.compile(
        r"(?<!\d)(\d[\d,]*(?:/\d[\d,]*)?\s*(?:개|건|종|행|셀|줄|지구|파일|문장|기록|오류)?)"
    )
    tokens: list[str] = []
    for item in items:
        for token in pattern.findall(str(item.get("text", ""))):
            normalized = token.replace(" ", "")
            if normalized not in tokens:
                tokens.append(normalized)
            if len(tokens) >= limit:
                return tokens
    return tokens


def _render_story_metrics(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
) -> None:
    groups = _group_decision_body(slide_spec.get("body", []))
    result_items = groups["결과"]
    metrics = _metric_tokens(result_items)
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage_reserved = 1.45 if slide_spec.get("stage") else 0.85
    top = 1.72
    bottom = canvas_h - stage_reserved
    left = 0.92
    main_w = (canvas_w - left * 2) * 0.46
    side_x = left + main_w + 0.36
    side_w = canvas_w - side_x - left
    main_h = bottom - top
    _add_shape(slide, MSO_SHAPE.RECTANGLE, (left, top, main_w, main_h), fill=style["primary_pale"])
    _add_textbox(
        slide,
        "핵심 확인값",
        (left + 0.30, top + 0.18, main_w - 0.60, 0.36),
        font=font,
        size=style["eyebrow_pt"],
        color=style["primary_strong"],
        bold=True,
    )
    _add_rule(slide, (left + 0.30, top + 0.58, main_w - 0.60, 0.04), style["primary"])
    metric_gap = 0.16
    metric_count = max(1, len(metrics))
    metric_w = (main_w - 0.60 - metric_gap * (metric_count - 1)) / metric_count
    for index, metric in enumerate(metrics or ["확인값"]):
        metric_x = left + 0.30 + index * (metric_w + metric_gap)
        _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            (metric_x, top + 0.78, metric_w, 1.08),
            fill=style["surface"],
            line=style["primary_soft"],
        )
        _add_textbox(
            slide,
            metric,
            (metric_x + 0.08, top + 0.88, metric_w - 0.16, 0.76),
            font=font,
            size=max(25.0, float(style["title_pt"])),
            color=style["primary_strong"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    _add_spread_copy(
        slide,
        result_items,
        (left + 0.30, top + 2.08, main_w - 0.60, main_h - 2.38),
        style,
        font,
        accent=style["primary"],
    )
    _add_labeled_panel(
        slide,
        "왜 했나",
        groups["결론"],
        (side_x, top, side_w, (bottom - top) * 0.48),
        style,
        font,
        accent=style["primary_strong"],
        fill=style["surface"],
        spread=True,
    )
    _add_labeled_panel(
        slide,
        "다음 방향",
        groups["남은 한계"],
        (side_x, top + (bottom - top) * 0.52, side_w, (bottom - top) * 0.48),
        style,
        font,
        accent=style["feedback"],
        fill=style["feedback_soft"],
        spread=True,
    )


def _process_positions(
    count: int,
    left: float,
    right: float,
    *,
    single_at_right: bool = False,
) -> list[float]:
    if count <= 0:
        return []
    if count == 1:
        return [right if single_at_right else (left + right) / 2]
    step = (right - left) / (count - 1)
    return [left + step * index for index in range(count)]


def _is_design_only(item: dict[str, Any]) -> bool:
    text = item.get("text", "")
    return any(marker in text for marker in ("설계만", "미구현", "구현 전", "후속 설계"))


def _add_horizontal_process_arrow(
    slide: Any,
    start: float,
    end: float,
    center_y: float,
    style: dict[str, Any],
) -> None:
    width = end - start
    if width <= 0.18:
        return
    _add_shape(
        slide,
        MSO_SHAPE.RIGHT_ARROW,
        (start, center_y - 0.07, width, 0.14),
        fill=style["primary_soft"],
        line=style["primary"],
    )


def _add_process_turn(
    slide: Any,
    start_x: float,
    start_y: float,
    end_x: float,
    end_y: float,
    icon_size: float,
    style: dict[str, Any],
) -> None:
    """첫째 행 끝에서 둘째 행 첫 노드로 이어지는 편집형 꺾임선을 그린다."""
    route_x = style["canvas_width"] - 0.72
    route_y = end_y - icon_size * 0.78
    if route_x > start_x + icon_size / 2:
        _add_rule(
            slide,
            (start_x + icon_size / 2, start_y - 0.012, route_x - start_x - icon_size / 2, 0.024),
            style["border"],
        )
    _add_rule(
        slide,
        (route_x, start_y, 0.024, route_y - start_y),
        style["border"],
    )
    if route_x > end_x:
        _add_rule(
            slide,
            (end_x, route_y, route_x - end_x, 0.024),
            style["border"],
        )
    arrow_height = max(0.20, end_y - route_y - icon_size / 2)
    _add_shape(
        slide,
        MSO_SHAPE.DOWN_ARROW,
        (end_x - 0.07, route_y, 0.14, arrow_height),
        fill=style["primary_soft"],
        line=style["primary"],
    )


def _render_process_ledger(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
) -> None:
    """긴 bridge 항목은 두 열 진행 목록으로 렌더해 빈 공간을 줄인다."""
    items = list(slide_spec.get("body", []))
    if not items:
        return
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage_reserved = 1.45 if slide_spec.get("stage") else 0.85
    top = 1.76
    bottom = canvas_h - stage_reserved
    left = 0.92
    gap = 0.28
    column_w = (canvas_w - left * 2 - gap) / 2
    rows = (len(items) + 1) // 2
    row_gap = 0.20
    row_h = (bottom - top - row_gap * max(0, rows - 1)) / rows
    silhouettes = (
        MSO_SHAPE.FLOWCHART_DOCUMENT,
        MSO_SHAPE.FLOWCHART_PROCESS,
        MSO_SHAPE.FLOWCHART_STORED_DATA,
        MSO_SHAPE.FLOWCHART_DECISION,
        MSO_SHAPE.GEAR_6,
        MSO_SHAPE.CHEVRON,
    )
    for index, item in enumerate(items):
        column = index % 2
        row = index // 2
        x = left + column * (column_w + gap)
        y = top + row * (row_h + row_gap) + (0.10 if column else 0.0)
        fill = style["feedback_soft"] if _is_design_only(item) else style["primary_pale"]
        ink = style["feedback"] if _is_design_only(item) else style["primary"]
        _add_shape(slide, MSO_SHAPE.RECTANGLE, (x, y, column_w, row_h), fill=fill)
        _add_rule(slide, (x, y, 0.055, row_h), ink)
        _add_shape(
            slide,
            silhouettes[index % len(silhouettes)],
            (x + 0.22, y + 0.22, 0.58, 0.58),
            fill=style["surface"],
            line=ink,
        )
        _add_textbox(
            slide,
            f"{index + 1:02d}",
            (x + 0.28, y + 0.38, 0.46, 0.24),
            font=font,
            size=style["caption_pt"],
            color=ink,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide,
            item["text"],
            (x + 0.98, y + 0.18, column_w - 1.20, row_h - 0.36),
            font=font,
            size=max(16.0, float(style["body_small_pt"])),
            color=style["feedback"] if _is_design_only(item) else style["text"],
            bold=bool(item.get("emphasis", False)) or _is_design_only(item),
        )


def _render_process(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
) -> None:
    """bridge body를 순서형 2행 PowerPoint 프로세스 맵으로 렌더한다."""
    items = list(slide_spec.get("body", []))
    design_only = [item for item in items if _is_design_only(item)]
    ordered = [item for item in items if not _is_design_only(item)] + design_only
    if not ordered:
        return
    if len(ordered) > 6 or max(len(item.get("text", "")) for item in ordered) >= 34:
        _render_process_ledger(slide, slide_spec, style, font)
        return

    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage_reserved = 1.45 if slide_spec.get("stage") else 0.85
    content_top = 1.78
    content_bottom = canvas_h - stage_reserved
    content_height = content_bottom - content_top
    icon_size = max(0.62, min(0.86, canvas_h * 0.078))

    if design_only and len(ordered) == 1:
        top_count = 0
    else:
        top_count = (len(ordered) + 1) // 2
    top_items = ordered[:top_count]
    bottom_items = ordered[top_count:]
    node_text_width = max(1.65, min(3.2, (canvas_w - 1.8) / max(2, len(top_items))))
    edge_clearance = max(1.45, node_text_width / 2 + 0.48)
    left_center = edge_clearance
    right_center = canvas_w - edge_clearance
    top_y = content_top + content_height * 0.22
    bottom_y = content_top + content_height * 0.67
    top_xs = _process_positions(len(top_items), left_center, right_center)
    bottom_xs = _process_positions(
        len(bottom_items),
        left_center,
        right_center,
        single_at_right=bool(design_only),
    )

    for start, end in zip(top_xs, top_xs[1:]):
        _add_horizontal_process_arrow(
            slide,
            start + icon_size / 2 + 0.10,
            end - icon_size / 2 - 0.10,
            top_y,
            style,
        )
    if top_xs and bottom_xs:
        _add_process_turn(
            slide,
            top_xs[-1],
            top_y,
            bottom_xs[0],
            bottom_y,
            icon_size,
            style,
        )
    for start, end in zip(bottom_xs, bottom_xs[1:]):
        _add_horizontal_process_arrow(
            slide,
            start + icon_size / 2 + 0.10,
            end - icon_size / 2 - 0.10,
            bottom_y,
            style,
        )

    silhouettes = (
        MSO_SHAPE.FLOWCHART_DOCUMENT,
        MSO_SHAPE.HEXAGON,
        MSO_SHAPE.FLOWCHART_PROCESS,
        MSO_SHAPE.FLOWCHART_DECISION,
        MSO_SHAPE.FLOWCHART_STORED_DATA,
        MSO_SHAPE.GEAR_6,
        MSO_SHAPE.CHEVRON,
    )
    nodes = list(zip(top_items, top_xs, [top_y] * len(top_items)))
    nodes.extend(zip(bottom_items, bottom_xs, [bottom_y] * len(bottom_items)))
    minimum = max(16.0, float(style["body_small_pt"]))
    for index, (item, center_x, center_y) in enumerate(nodes):
        feedback = _is_design_only(item)
        ink = style["feedback"] if feedback else style["primary"]
        soft = style["feedback_soft"] if feedback else style["primary_soft"]
        _add_shape(
            slide,
            silhouettes[index % len(silhouettes)],
            (center_x - icon_size / 2, center_y - icon_size / 2, icon_size, icon_size),
            fill=soft,
            line=ink,
        )
        _add_textbox(
            slide,
            f"{index + 1:02d}",
            (center_x - 0.28, center_y - icon_size * 0.20, 0.56, 0.28),
            font=font,
            size=style["caption_pt"],
            color=ink,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if feedback:
            _add_textbox(
                slide,
                "설계만",
                (center_x - 0.58, center_y - icon_size * 0.95, 1.16, 0.28),
                font=font,
                size=style["caption_pt"],
                color=style["feedback"],
                bold=True,
                align=PP_ALIGN.CENTER,
            )
        _add_textbox(
            slide,
            item["text"],
            (center_x - node_text_width / 2, center_y + icon_size * 0.64, node_text_width, 0.82),
            font=font,
            size=minimum,
            color=style["feedback"] if feedback else style["text"],
            bold=bool(item.get("emphasis", False)) or feedback,
            align=PP_ALIGN.CENTER,
        )


def _render_process_stair(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
) -> None:
    items = list(slide_spec.get("body", []))
    reason_items = [item for item in items if item.get("kind") == "범위"]
    steps = [item for item in items if item.get("kind") != "범위"]
    if not steps:
        return
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage_reserved = 1.45 if slide_spec.get("stage") else 0.85
    top = 1.76
    bottom = canvas_h - stage_reserved
    reason_w = max(2.65, canvas_w * 0.22)
    reason_item = reason_items[0] if reason_items else {
        "text": "읽는 기준 — 입력 범위를 먼저 고정하고 확인 순서와 다음 입력을 구분함",
        "level": 1,
    }
    checkpoint_items = [
        reason_item,
        {"text": f"시작 — {_process_checkpoint(steps[0]['text'])}", "level": 2},
        {"text": f"중간 — {_process_checkpoint(steps[len(steps) // 2]['text'])}", "level": 2},
        {"text": f"다음 — {_process_checkpoint(steps[-1]['text'])}", "level": 2},
    ]
    _add_labeled_panel(
        slide,
        "왜 이 순서인가",
        checkpoint_items,
        (0.92, top, reason_w, bottom - top),
        style,
        font,
        accent=style["primary_strong"],
        fill=style["primary_pale"],
        spread=True,
    )
    grid_left = 0.92 + reason_w + 0.34
    grid_w = canvas_w - grid_left - 0.92
    rows = 2
    columns = max(1, (len(steps) + 1) // rows)
    card_w = (grid_w - 0.24 * max(0, columns - 1)) / columns
    card_h = (bottom - top - 0.34) / rows
    silhouettes = (
        MSO_SHAPE.FLOWCHART_DOCUMENT,
        MSO_SHAPE.HEXAGON,
        MSO_SHAPE.FLOWCHART_PROCESS,
        MSO_SHAPE.FLOWCHART_DECISION,
        MSO_SHAPE.FLOWCHART_STORED_DATA,
        MSO_SHAPE.GEAR_6,
    )
    for index, item in enumerate(steps):
        row = index // columns
        col = index % columns
        x = grid_left + col * (card_w + 0.24)
        y = top + row * (card_h + 0.34) + col * 0.04
        design_only = _is_design_only(item)
        ink = style["feedback"] if design_only else style["primary"]
        fill = style["feedback_soft"] if design_only else style["surface"]
        _add_shape(slide, MSO_SHAPE.RECTANGLE, (x, y, card_w, card_h), fill=fill, line=style["border"])
        _add_shape(
            slide,
            silhouettes[index % len(silhouettes)],
            (x + 0.22, y + 0.24, 0.62, 0.62),
            fill=style["primary_soft"] if not design_only else style["feedback_soft"],
            line=ink,
        )
        _add_textbox(
            slide,
            f"{index + 1:02d}",
            (x + 0.26, y + 0.42, 0.54, 0.24),
            font=font,
            size=style["caption_pt"],
            color=ink,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if design_only:
            _add_textbox(
                slide,
                "설계만",
                (x + card_w - 0.92, y + 0.16, 0.68, 0.28),
                font=font,
                size=style["caption_pt"],
                color=style["feedback"],
                bold=True,
                align=PP_ALIGN.RIGHT,
            )
        _add_textbox(
            slide,
            item["text"],
            (x + 0.22, y + 0.70, card_w - 0.44, card_h - 0.84),
            font=font,
            size=max(16.0, float(style["body_small_pt"])),
            color=ink if design_only else style["text"],
            bold=design_only or bool(item.get("emphasis", False)),
        )


def _process_checkpoint(text: str) -> str:
    """절차 카드의 기존 문장에서 왼쪽 읽기 레일에 쓸 짧은 표지를 고른다."""
    if " — " in text:
        return text.split(" — ", 1)[0]
    labels = (
        ("목적별", "비교 기능 후속"),
        ("누락 사유", "누락 사유 인계"),
        ("첨부 목록", "첨부 목록 대조"),
        ("원본 파일", "원본 보존"),
        ("형식과 열림", "파일 열림 확인"),
        ("도면 파일", "도면 분리 기록"),
        ("전국 지구 목록", "대상 목록 고정"),
        ("기존 구조", "기존 구조 폐기"),
        ("항목 이름", "기준 파일 고정"),
        ("실제 대상", "대상·기록 분리"),
        ("고시 문서", "문서·사건 분리"),
        ("기준일", "유효 계획판 선택"),
        ("조건 충족", "미확인 값 분리"),
    )
    for marker, label in labels:
        if marker in text:
            return label
    clause = text.split(",", 1)[0].strip()
    return clause if len(clause) <= 18 else f"{clause[:17]}…"


def _render_overview_asym(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
) -> None:
    """원문 수집과 의미 비교를 서로 다른 폭의 두 작업축으로 배치한다."""
    items = list(slide_spec.get("body", []))
    if not items:
        return
    canvas_w = style["canvas_width"]
    intro = items[0]
    _add_rule(slide, (0.92, 1.77, 0.08, 0.48), style["primary"])
    _add_textbox(
        slide,
        str(intro["text"]),
        (1.18, 1.72, canvas_w - 2.10, 0.58),
        font=font,
        size=max(14.0, float(style["caption_pt"])),
        color=style["primary_strong"],
        bold=bool(intro.get("emphasis", False)),
    )

    footnote = items[-1] if str(items[-1].get("text", "")).startswith("내주.") else None
    content_items = items[1:-1] if footnote else items[1:]
    split_at = next(
        (
            index
            for index, item in enumerate(content_items)
            if str(item.get("text", "")).startswith("조·항의 뜻을 비교하는 기준 —")
        ),
        max(1, len(content_items) // 2),
    )
    collection_items = content_items[:split_at]
    meaning_items = content_items[split_at:]
    if not collection_items or not meaning_items:
        _render_process_ledger(slide, slide_spec, style, font)
        return

    left_x = 0.92
    content_w = canvas_w - 1.84
    gap = 0.54
    left_w = content_w * 0.59
    right_w = content_w - left_w - gap
    right_x = left_x + left_w + gap
    left_y = 2.46
    right_y = 2.68
    panel_bottom = style["canvas_height"] - (1.02 if footnote else 0.78)
    left_h = panel_bottom - left_y
    right_h = panel_bottom - right_y

    def add_axis(
        axis_items: list[dict[str, Any]],
        box: tuple[float, float, float, float],
        *,
        accent: str,
        fill: str,
        compact: bool,
    ) -> None:
        x, y, width, height = box
        heading = axis_items[0]
        steps = axis_items[1:]
        header_h = 1.04
        _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            (x, y, width, header_h),
            fill=accent,
            line=accent,
        )
        _add_shape(
            slide,
            MSO_SHAPE.FLOWCHART_STORED_DATA if not compact else MSO_SHAPE.HEXAGON,
            (x + 0.22, y + 0.22, 0.58, 0.58),
            fill=style["surface"],
            line=style["surface"],
        )
        _add_label_content_textbox(
            slide,
            str(heading["text"]),
            (x + 1.02, y + 0.08, width - 1.28, 0.86),
            font=font,
            label_size=max(17.0, float(style["body_small_pt"])),
            content_size=max(13.5, float(style["caption_pt"])),
            label_color=style["surface"],
            content_color=style["primary_soft"],
        )
        if not steps:
            return
        steps_y = y + header_h + 0.15
        steps_h = height - header_h - 0.15
        step_gap = 0.08
        step_h = (steps_h - step_gap * (len(steps) - 1)) / len(steps)
        for index, item in enumerate(steps):
            step_y = steps_y + index * (step_h + step_gap)
            is_limit = item.get("kind") == "남은 한계"
            band_fill = style["feedback_soft"] if is_limit else fill
            band_accent = style["feedback"] if is_limit else accent
            _add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                (x, step_y, width, step_h),
                fill=band_fill,
                line=style["border"],
            )
            _add_rule(slide, (x, step_y, 0.055, step_h), band_accent)
            _add_shape(
                slide,
                MSO_SHAPE.OVAL,
                (x + 0.18, step_y + max(0.10, (step_h - 0.48) / 2), 0.48, 0.48),
                fill=style["surface"],
                line=band_accent,
            )
            _add_textbox(
                slide,
                f"{index + 1:02d}",
                (x + 0.21, step_y + max(0.21, (step_h - 0.18) / 2), 0.42, 0.18),
                font=font,
                size=max(10.5, float(style["caption_pt"]) - 1.5),
                color=band_accent,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            _add_label_content_textbox(
                slide,
                str(item["text"]),
                (x + 0.84, step_y + 0.06, width - 1.06, step_h - 0.12),
                font=font,
                label_size=max(14.0, float(style["caption_pt"])),
                content_size=max(12.0, float(style["caption_pt"]) - (0.5 if compact else 0.0)),
                label_color=band_accent,
                content_color=style["text"],
            )

    add_axis(
        collection_items,
        (left_x, left_y, left_w, left_h),
        accent=style["primary_strong"],
        fill=style["primary_pale"],
        compact=False,
    )
    add_axis(
        meaning_items,
        (right_x, right_y, right_w, right_h),
        accent=style["primary"],
        fill=style["surface_subtle"],
        compact=True,
    )
    _add_textbox(
        slide,
        "정본 입력",
        (left_x + left_w + 0.03, 5.55, gap - 0.06, 0.30),
        font=font,
        size=max(10.0, float(style["caption_pt"]) - 2.0),
        color=style["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_shape(
        slide,
        MSO_SHAPE.RIGHT_ARROW,
        (left_x + left_w + 0.08, 5.92, gap - 0.16, 0.26),
        fill=style["primary_soft"],
        line=style["primary"],
    )
    if footnote:
        _add_rule(slide, (0.92, panel_bottom + 0.16, 0.10, 0.32), style["feedback"])
        _add_textbox(
            slide,
            str(footnote["text"]),
            (1.18, panel_bottom + 0.10, canvas_w - 2.10, 0.44),
            font=font,
            size=max(10.5, float(style["caption_pt"]) - 1.0),
            color=style["muted"],
            bold=True,
        )


def _add_evidence_strip(
    slide: Any,
    items: list[dict[str, Any]],
    box: tuple[float, float, float, float],
    style: dict[str, Any],
    font: str,
) -> None:
    x, y, width, height = box
    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        box,
        fill=style["primary_pale"],
    )
    feedback = any(item.get("kind") == "남은 한계" for item in items)
    accent = style["feedback"] if feedback else style["primary"]
    _add_rule(slide, (x, y, 0.055, height), accent)
    label = next((item.get("kind") for item in items if item.get("kind")), "해설")
    _add_textbox(
        slide,
        str(label),
        (x + 0.18, y + 0.10, 1.15, height - 0.20),
        font=font,
        size=style["caption_pt"],
        color=accent,
        bold=True,
    )
    shape = slide.shapes.add_textbox(
        Inches(x + 1.35),
        Inches(y + 0.08),
        Inches(width - 1.53),
        Inches(height - 0.16),
    )
    frame = shape.text_frame
    _reset_text_frame(frame, vertical_anchor=MSO_ANCHOR.MIDDLE)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item["text"]
        paragraph.space_after = Pt(3)
        paragraph.font.name = font
        paragraph.font.size = Pt(max(16.0, float(style["body_small_pt"])))
        paragraph.font.bold = bool(item.get("emphasis", False))
        paragraph.font.underline = bool(item.get("emphasis", False))
        paragraph.font.color.rgb = _hex(style["text"])


def _render_visual_focus(
    slide: Any,
    slide_spec: dict[str, Any],
    spec_path: Path,
    style: dict[str, Any],
    font: str,
) -> None:
    """두 evidence strip 아래에 bridge image를 가능한 크게 유지한다."""
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage = slide_spec.get("stage")
    visual = slide_spec["visuals"][0]
    caption = visual.get("caption")
    caption_height = 0.5 if caption else 0.0
    side_margin = 0.9
    content_width = canvas_w - side_margin * 2
    items = list(slide_spec.get("body", []))
    strip_gap = 0.22
    strip_y = 1.70
    strip_height = max(0.86, min(1.08, canvas_h * 0.095)) if items else 0.0
    split = max(1, (len(items) + 1) // 2)
    groups = (items[:split], items[split:]) if len(items) > 1 else (items, [])
    if groups[1]:
        strip_width = (content_width - strip_gap) / 2
        _add_evidence_strip(
            slide,
            groups[0],
            (side_margin, strip_y, strip_width, strip_height),
            style,
            font,
        )
        _add_evidence_strip(
            slide,
            groups[1],
            (side_margin + strip_width + strip_gap, strip_y, strip_width, strip_height),
            style,
            font,
        )
    elif groups[0]:
        _add_evidence_strip(
            slide,
            groups[0],
            (side_margin, strip_y, content_width, strip_height),
            style,
            font,
        )

    stage_reserved = 1.45 if stage else 0.85
    visual_y = strip_y + strip_height + (0.14 if items else 0.0)
    visual_height = max(1.0, canvas_h - visual_y - stage_reserved)
    visual_box = (
        side_margin,
        visual_y,
        content_width,
        visual_height - caption_height,
    )
    asset = _resolve_asset(visual["path"], spec_path)
    if visual["kind"] == "image":
        _add_image(
            slide,
            asset,
            visual_box,
            crop_bottom=float(visual.get("crop_bottom", 0.0)),
        )
    elif visual["kind"] == "table":
        _add_table(slide, asset, visual_box, style, font)
    if caption:
        _add_textbox(
            slide,
            caption,
            (visual_box[0], visual_y + visual_height - caption_height, visual_box[2], caption_height),
            font=font,
            size=style["caption_pt"],
            color=style["muted"],
        )


def _render_visual_rail(
    slide: Any,
    slide_spec: dict[str, Any],
    spec_path: Path,
    style: dict[str, Any],
    font: str,
) -> None:
    if not slide_spec.get("visuals"):
        return
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage_reserved = 1.45 if slide_spec.get("stage") else 0.85
    top = 1.72
    bottom = canvas_h - stage_reserved
    left = 0.92
    content_w = canvas_w - left * 2
    rail_w = content_w * (0.48 if canvas_w < 15.0 else 0.18)
    gap = 0.34
    visual_x = left + rail_w + gap
    visual_w = content_w - rail_w - gap
    rail_items = list(slide_spec.get("body", []))
    _add_labeled_panel(
        slide,
        "온톨로지 구성",
        rail_items,
        (left, top, rail_w, bottom - top),
        style,
        font,
        accent=style["primary_strong"],
        fill=style["primary_pale"],
        watermark=False,
    )
    visual = slide_spec["visuals"][0]
    caption = visual.get("caption")
    caption_h = 0.46 if caption else 0.0
    visual_box = (visual_x, top, visual_w, bottom - top - caption_h)
    asset = _resolve_asset(visual["path"], spec_path)
    if visual["kind"] == "image":
        _add_image(
            slide,
            asset,
            visual_box,
            crop_bottom=float(visual.get("crop_bottom", 0.0)),
        )
    elif visual["kind"] == "table":
        _add_table(slide, asset, visual_box, style, font)
    if caption:
        _add_textbox(
            slide,
            caption,
            (visual_x, bottom - caption_h, visual_w, caption_h),
            font=font,
            size=style["caption_pt"],
            color=style["muted"],
        )


def _render_summary_asym(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
) -> None:
    groups = _group_decision_body(slide_spec.get("body", []))
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage_reserved = 1.45 if slide_spec.get("stage") else 0.85
    top = 1.72
    bottom = canvas_h - stage_reserved
    left = 0.92
    gap = 0.34
    flow_w = (canvas_w - left * 2 - gap) * 0.58
    next_w = canvas_w - left * 2 - gap - flow_w
    _add_labeled_panel(
        slide,
        "확인한 흐름",
        groups["결과"] + groups["결론"],
        (left, top, flow_w, bottom - top),
        style,
        font,
        accent=style["primary_strong"],
        fill=style["primary_pale"],
        dense=False,
        spread=True,
    )
    _add_labeled_panel(
        slide,
        "다음 방향",
        groups["남은 한계"],
        (left + flow_w + gap, top + 0.34, next_w, bottom - top - 0.34),
        style,
        font,
        accent=style["feedback"],
        fill=style["feedback_soft"],
        spread=True,
    )


def _render_glossary(
    slide: Any,
    slide_spec: dict[str, Any],
    style: dict[str, Any],
    font: str,
) -> None:
    items = list(slide_spec.get("body", []))
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage_reserved = 1.45 if slide_spec.get("stage") else 0.85
    top = 1.72
    bottom = canvas_h - stage_reserved
    left = 0.92
    content_w = canvas_w - left * 2
    content_h = bottom - top
    intro_w = content_w * 0.28
    gap = 0.28
    grid_x = left + intro_w + gap
    grid_w = content_w - intro_w - gap

    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        (left, top + 0.24, intro_w, content_h - 0.24),
        fill=style["primary_pale"],
        line=style["border"],
    )
    _add_shape(
        slide,
        MSO_SHAPE.OVAL,
        (left + 0.24, top + 0.50, 0.72, 0.72),
        fill=style["primary"],
        line=style["primary"],
    )
    _add_textbox(
        slide,
        "Aa",
        (left + 0.24, top + 0.64, 0.72, 0.28),
        font=font,
        size=style["caption_pt"],
        color=style["surface"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        "프로젝트 용어 정의",
        (left + 0.24, top + 1.34, intro_w - 0.48, 0.72),
        font=font,
        size=style["eyebrow_pt"],
        color=style["primary_strong"],
        bold=True,
    )
    _add_textbox(
        slide,
        "분석용 문서, 조문 정본, 문언후보, 정의환경, 판단 기록, 온톨로지는 이 장에 적은 대상과 범위를 기준으로 사용한다.",
        (left + 0.24, top + 2.18, intro_w - 0.48, 2.10),
        font=font,
        size=style["body_small_pt"],
        color=style["text"],
    )
    _add_textbox(
        slide,
        "이전 화면으로 돌아가기",
        (left + 0.24, bottom - 0.82, intro_w - 0.48, 0.38),
        font=font,
        size=style["caption_pt"],
        color=style["feedback"],
        bold=True,
        underline=True,
    )

    if not items:
        return
    col_gap = 0.22
    cols = 2
    rows = (len(items) + cols - 1) // cols
    cell_w = (grid_w - col_gap) / cols
    row_gap = 0.14
    cell_h = max(0.82, (content_h - row_gap * (rows - 1)) / rows)
    for index, item in enumerate(items):
        col = index // rows
        row = index % rows
        x = grid_x + col * (cell_w + col_gap)
        y = top + row * (cell_h + row_gap)
        emphasized = bool(item.get("emphasis", False)) or int(item.get("level", 1)) == 1
        fill = style["surface"] if not emphasized else style["surface_subtle"]
        _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            (x, y, cell_w, cell_h),
            fill=fill,
            line=style["border"],
        )
        _add_textbox(
            slide,
            f"{index + 1:02d}",
            (x + 0.16, y + 0.14, 0.44, 0.28),
            font=font,
            size=style["caption_pt"],
            color=style["primary"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide,
            str(item["text"]),
            (x + 0.70, y + 0.11, cell_w - 0.88, cell_h - 0.22),
            font=font,
            size=max(14.0, float(style["body_small_pt"])),
            color=style["primary_strong"] if emphasized else style["text"],
            bold=emphasized,
        )


def _render_bridge_visual_content(
    slide: Any,
    slide_spec: dict[str, Any],
    spec_path: Path,
    style: dict[str, Any],
    font: str,
    *,
    title_y: float,
    title_height: float,
) -> None:
    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage = slide_spec.get("stage")
    visual = slide_spec["visuals"][0]
    caption = visual.get("caption")
    caption_height = 0.5 if caption else 0.0
    side_margin = 0.9
    content_width = canvas_w - side_margin * 2
    stage_reserved = 1.45 if stage else 0.85
    body_y = title_y + title_height + 0.25
    body_height = max(0.8, min(1.05, canvas_h * 0.085))
    gap = 0.22
    visual_y = body_y + body_height + gap
    visual_height = max(1.0, canvas_h - visual_y - stage_reserved)
    visual_box = (
        side_margin,
        visual_y,
        content_width,
        visual_height - caption_height,
    )

    _add_body(
        slide,
        slide_spec.get("body", []),
        (side_margin, body_y, content_width, body_height),
        style,
        font,
    )
    asset = _resolve_asset(visual["path"], spec_path)
    if visual["kind"] == "image":
        _add_image(
            slide,
            asset,
            visual_box,
            crop_bottom=float(visual.get("crop_bottom", 0.0)),
        )
    elif visual["kind"] == "table":
        _add_table(slide, asset, visual_box, style, font)
    if caption:
        _add_textbox(
            slide,
            caption,
            (visual_box[0], visual_y + visual_height - caption_height, visual_box[2], caption_height),
            font=font,
            size=style["caption_pt"],
            color=style["muted"],
        )


def _render_slide(
    prs: Presentation,
    slide_spec: dict[str, Any],
    spec_path: Path,
    style: dict[str, Any],
    font: str,
    *,
    as_of: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cover = slide_spec["role"] == "cover"
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _hex(style["primary"] if cover else style["surface"])

    canvas_w = style["canvas_width"]
    canvas_h = style["canvas_height"]
    stage = slide_spec.get("stage")
    layout = slide_spec.get("layout", "default")
    if cover and layout == "cover-asym":
        _render_cover_asym(slide, slide_spec, style, font, as_of=as_of)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return
    if cover:
        title_x = 0.9
        title_width = canvas_w - 1.8
    else:
        title_x, title_width = _add_noncover_chrome(slide, slide_spec, style, font)
    title_y = 1.55 if cover else 0.55
    title_height = 1.8 if cover else 0.8
    title_size = style["cover_title_pt"] if cover else style["title_pt"]
    _add_textbox(
        slide,
        slide_spec["title"],
        (title_x, title_y, title_width, title_height),
        font=font,
        size=title_size,
        color=style["surface"] if cover else style["text"],
        bold=True,
        align=PP_ALIGN.CENTER if cover else PP_ALIGN.LEFT,
    )
    if cover:
        _add_cover_chrome(slide, as_of=as_of, style=style, font=font)

    visuals = slide_spec.get("visuals", [])
    if layout == "decision-columns":
        _render_decision_columns(slide, slide_spec, style, font)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return
    if layout in ("story-left", "story-right"):
        _render_story_asym(slide, slide_spec, style, font, layout=layout)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return
    if layout == "story-metrics":
        _render_story_metrics(slide, slide_spec, style, font)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return
    if layout == "process":
        _render_process(slide, slide_spec, style, font)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return
    if layout == "process-stair":
        _render_process_stair(slide, slide_spec, style, font)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return
    if layout == "overview-asym":
        _render_overview_asym(slide, slide_spec, style, font)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return
    if layout == "visual-focus" and len(visuals) == 1:
        _render_visual_focus(slide, slide_spec, spec_path, style, font)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return
    if layout == "visual-rail" and len(visuals) == 1:
        _render_visual_rail(slide, slide_spec, spec_path, style, font)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return
    if layout == "summary-asym":
        _render_summary_asym(slide, slide_spec, style, font)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return
    if layout == "glossary":
        _render_glossary(slide, slide_spec, style, font)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return
    if slide_spec["role"] == "bridge" and len(visuals) == 1:
        _render_bridge_visual_content(
            slide,
            slide_spec,
            spec_path,
            style,
            font,
            title_y=title_y,
            title_height=title_height,
        )
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if stage:
            _add_stage_rail(slide, int(stage), style, font)
        return

    body_y = 4.1 if cover else 2.05
    body_height = max(1.0, canvas_h - body_y - (1.45 if stage else 0.85))
    body_width = canvas_w - 1.8 if not visuals else (canvas_w - 2.1) * 0.48
    _add_body(
        slide,
        slide_spec.get("body", []),
        (0.9, body_y, body_width, body_height),
        style,
        font,
        color=style["surface"] if cover else style["text"],
    )

    if visuals:
        visual = visuals[0]
        asset = _resolve_asset(visual["path"], spec_path)
        caption = visual.get("caption")
        caption_height = 0.5 if caption else 0.0
        visual_box = (
            canvas_w * 0.53,
            body_y,
            canvas_w * 0.42,
            body_height - caption_height,
        )
        if visual["kind"] == "image":
            _add_image(
                slide,
                asset,
                visual_box,
                crop_bottom=float(visual.get("crop_bottom", 0.0)),
            )
        elif visual["kind"] == "table":
            _add_table(slide, asset, visual_box, style, font)
        if caption:
            _add_textbox(
                slide,
                caption,
                (visual_box[0], body_y + body_height - caption_height, visual_box[2], caption_height),
                font=font,
                size=style["caption_pt"],
                color=style["muted"],
            )

    notes = slide_spec.get("notes")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    if stage:
        _add_stage_rail(slide, int(stage), style, font)


def _wire_internal_links(prs: Presentation) -> None:
    """본문의 용어 링크를 부록에 연결하고 부록에는 이전 화면 복귀를 건다."""
    glossary_slide = None
    for slide in prs.slides:
        if any(
            shape.has_text_frame and shape.text.strip() == "용어 설명"
            for shape in slide.shapes
        ):
            glossary_slide = slide
            break
    if glossary_slide is None:
        return

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            label = shape.text.strip()
            if label == "용어 설명 ↗":
                shape.click_action.target_slide = glossary_slide
            elif label == "이전 화면으로 돌아가기":
                hyperlink = shape._element.nvSpPr.cNvPr.get_or_add_hlinkClick()
                hyperlink.action = "ppaction://hlinkshowjump?jump=lastslideviewed"


def _render_previews(pptx_path: Path, out_dir: Path, slide_count: int) -> tuple[str | None, list[str], list[str]]:
    notes: list[str] = []
    soffice = shutil.which("soffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        missing = [name for name, value in (("soffice", soffice), ("pdftoppm", pdftoppm)) if not value]
        notes.append(f"미리보기 도구 미설치: {', '.join(missing)}")
        return None, [], notes

    pdf_path = out_dir / "deck.pdf"
    previews_dir = out_dir / "previews"
    previews_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="ppt-render-") as temp_dir:
        profile_uri = Path(temp_dir).resolve().as_uri()
        converted = subprocess.run(
            [
                soffice,
                "--headless",
                f"-env:UserInstallation={profile_uri}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(out_dir),
                str(pptx_path),
            ],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )
        if converted.returncode != 0 or not pdf_path.is_file():
            notes.append(f"PDF 변환 실패: {converted.stderr.strip() or converted.stdout.strip()}")
            return None, [], notes
        prefix = Path(temp_dir) / "slide"
        rendered = subprocess.run(
            [pdftoppm, "-png", "-r", "80", str(pdf_path), str(prefix)],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )
        if rendered.returncode != 0:
            notes.append(f"PNG 변환 실패: {rendered.stderr.strip()}")
            return str(pdf_path), [], notes
        generated = sorted(Path(temp_dir).glob("slide-*.png"))
        previews: list[str] = []
        for number, source in enumerate(generated, start=1):
            target = previews_dir / f"s{number:02d}.png"
            shutil.copyfile(source, target)
            previews.append(str(target))
        if len(previews) != slide_count:
            notes.append(f"미리보기 수 불일치: {len(previews)}/{slide_count}")
        return str(pdf_path), previews, notes


def build(
    spec_path: Path | str,
    out_dir: Path | str,
    *,
    render_preview: bool = True,
    visual_check_slides: Iterable[int] = (),
    visual_check_note: str | None = None,
) -> dict[str, Any]:
    """덱 스펙을 PPTX로 렌더하고 빌드 자기보고를 반환한다."""
    spec_file = Path(spec_path).resolve()
    output = Path(out_dir).resolve()
    spec = json.loads(spec_file.read_text(encoding="utf-8"))
    schema = json.loads(DECK_SCHEMA.read_text(encoding="utf-8"))
    jsonschema.validate(spec, schema)

    style = apply_deck_overrides(load_style(spec["profile"]), spec)
    font = select_font(style["font"])
    prs = Presentation()
    prs.slide_width = Inches(style["canvas_width"])
    prs.slide_height = Inches(style["canvas_height"])
    for slide_spec in spec["slides"]:
        _render_slide(
            prs,
            slide_spec,
            spec_file,
            style,
            font["used"],
            as_of=spec["as_of"],
        )
    _wire_internal_links(prs)

    output.mkdir(parents=True, exist_ok=True)
    pptx_path = output / "deck.pptx"
    prs.save(pptx_path)

    notes: list[str] = []
    if font["fallback"]:
        notes.append(font["reason"])
    if render_preview:
        pdf, previews, preview_notes = _render_previews(pptx_path, output, len(spec["slides"]))
        notes.extend(preview_notes)
    else:
        pdf, previews = None, []
        notes.append("미리보기 렌더 비활성화")

    seen = sorted(set(int(value) for value in visual_check_slides))
    if seen and not (visual_check_note and visual_check_note.strip()):
        raise ValueError("육안확인 슬라이드를 지정하면 확인 대상·결과 메모가 필요함")
    report = {
        "deck_id": spec["deck_id"],
        "spec_sha256": sha256_file(spec_file),
        "source_sha256": spec["source_sha256"],
        "profile": spec["profile"],
        "slide_count": len(spec["slides"]),
        "font": font,
        "artifacts": {
            "pptx": str(pptx_path),
            "pdf": pdf,
            "previews": previews,
        },
        "visual_check": {
            "done": bool(seen),
            "slides_seen": seen,
            "note": visual_check_note.strip() if visual_check_note else None,
        },
        "notes": notes,
    }
    build_schema = json.loads(BUILD_SCHEMA.read_text(encoding="utf-8"))
    jsonschema.validate(report, build_schema)
    (output / "_build_report.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return report


def _parse_slides(value: str | None) -> list[int]:
    if not value:
        return []
    return [int(part) for part in value.split(",") if part.strip()]


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("spec", type=Path, help="02_deck.json 경로")
    parser.add_argument("--out", type=Path, help="build 디렉터리")
    parser.add_argument("--no-preview", action="store_true")
    parser.add_argument("--visual-check", help="육안확인한 슬라이드 번호. 예: 1,2")
    parser.add_argument("--visual-check-note", help="확인한 미리보기와 잘림·겹침 판정")
    args = parser.parse_args()

    output = args.out or args.spec.parent / "build"
    output_root = (ROOT / "output/ppt").resolve()
    if not args.spec.resolve().is_relative_to(output_root) or not output.resolve().is_relative_to(
        output_root
    ):
        parser.error("spec과 out은 output/ppt/<deck-id> 아래여야 함")
    report = build(
        args.spec,
        output,
        render_preview=not args.no_preview,
        visual_check_slides=_parse_slides(args.visual_check),
        visual_check_note=args.visual_check_note,
    )
    print(json.dumps(report, ensure_ascii=False, indent=2))
    complete = (
        report["artifacts"]["pdf"] is not None
        and len(report["artifacts"]["previews"]) == report["slide_count"]
        and report["visual_check"]["done"]
    )
    return 0 if complete else 2


if __name__ == "__main__":
    raise SystemExit(main())
